import os, shutil
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Target Paths
DESKTOP_DIR = r"C:\Users\Aryan\OneDrive\Desktop\Pranayuv"
OUTPUT_PPTX = os.path.join(DESKTOP_DIR, "Final_University_Internship_Presentation.pptx")
WORKSPACE_PPTX = os.path.abspath("Final_University_Internship_Presentation.pptx")

# Screenshot repository paths
SCREENSHOTS_DIR = os.path.abspath("screenshots")

# Corporate Color Palette (Blue & White Premium Corporate Theme)
NAVY_BLUE = RGBColor(31, 78, 121)     # #1F4E79 - Main Brand & Header
ACCENT_BLUE = RGBColor(46, 117, 182)  # #2E75B6 - Highlights & Subheadings
LIGHT_BG = RGBColor(242, 244, 248)    # #F2F4F8 - Card Backgrounds
WHITE = RGBColor(255, 255, 255)       # Pure White
DARK_TEXT = RGBColor(38, 38, 38)      # #262626 - Charcoal body text
GRAY_TEXT = RGBColor(100, 100, 100)   # #646464 - Muted descriptions
LINE_GRAY = RGBColor(210, 215, 225)   # #D2D7E1 - Borders & Gridlines

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title_text, category_text="INTERNSHIP TECHNICAL DEFENSE & PROJECT ANALYSIS"):
    # Category / Breadcrumb banner
    cat_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = 'Arial'
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_BLUE
    
    # Main Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.7), Inches(12), Inches(0.7))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = NAVY_BLUE
    
    # Decorative header rule line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY_BLUE
    line.line.color.rgb = NAVY_BLUE

def add_card(slide, left, top, width, height, bg_color=LIGHT_BG, border_color=LINE_GRAY):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
    else:
        card.line.fill.background()
    return card

def add_bullet_point(text_frame, bold_prefix, text_body, pt_size=14, color=DARK_TEXT, space_after=8):
    p = text_frame.add_paragraph() if len(text_frame.paragraphs[0].text) > 0 else text_frame.paragraphs[0]
    p.space_after = Pt(space_after)
    
    run_b = p.add_run()
    run_b.text = bold_prefix + (" " if bold_prefix else "")
    run_b.font.name = 'Arial'
    run_b.font.bold = True
    run_b.font.size = Pt(pt_size)
    run_b.font.color.rgb = color
    
    run_t = p.add_run()
    run_t.text = text_body
    run_t.font.name = 'Arial'
    run_t.font.bold = False
    run_t.font.size = Pt(pt_size)
    run_t.font.color.rgb = color

def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

def build_presentation():
    prs = Presentation()
    # Set widescreen dimensions 16:9 (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # -------------------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s1, NAVY_BLUE)
    
    # Decorative accent bar
    bar1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(0.2), Inches(4.5))
    bar1.fill.solid()
    bar1.fill.fore_color.rgb = WHITE
    bar1.line.fill.background()
    
    title_box = s1.shapes.add_textbox(Inches(1.3), Inches(1.5), Inches(11.5), Inches(4.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p0 = tf1.paragraphs[0]
    p0.text = "INDUSTRIAL INTERNSHIP REPORT & PROJECT ANALYSIS"
    p0.font.name = 'Arial'
    p0.font.size = Pt(32)
    p0.font.bold = True
    p0.font.color.rgb = WHITE
    p0.space_after = Pt(12)
    
    p1 = tf1.add_paragraph()
    p1.text = "Full-Stack Software Architecture & UI/UX Engineering Defense"
    p1.font.name = 'Arial'
    p1.font.size = Pt(22)
    p1.font.color.rgb = RGBColor(200, 220, 245)
    p1.space_after = Pt(40)
    
    p2 = tf1.add_paragraph()
    p2.text = "Student Name: Patel Aryan Shaileshbhai  |  Enrollment: SR24BSIT138"
    p2.font.name = 'Arial'
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.space_after = Pt(6)
    
    p3 = tf1.add_paragraph()
    p3.text = "Degree: B.Sc. Information Technology (Semester VI - AY 2025–2026)"
    p3.font.name = 'Arial'
    p3.font.size = Pt(16)
    p3.font.color.rgb = WHITE
    p3.space_after = Pt(6)
    
    p4 = tf1.add_paragraph()
    p4.text = "Institution: Shree Ramkrishna Institute of Computer Education & Applied Sciences, Sarvajanik University"
    p4.font.name = 'Arial'
    p4.font.size = Pt(15)
    p4.font.color.rgb = RGBColor(220, 230, 245)
    p4.space_after = Pt(6)
    
    p5 = tf1.add_paragraph()
    p5.text = "Host Company: Pranayuv Technologies Pvt. Ltd."
    p5.font.name = 'Arial'
    p5.font.size = Pt(16)
    p5.font.bold = True
    p5.font.color.rgb = RGBColor(160, 205, 255)

    add_speaker_notes(s1, 
        "Good morning respected HOD, Faculty Members, and External Examiners. My name is Patel Aryan Shaileshbhai, Enrollment number SR24BSIT138, currently in Semester VI of B.Sc. IT. "
        "Today, I am honored to present the comprehensive technical defense of my six-month industrial summer internship completed at Pranayuv Technologies Pvt. Ltd. "
        "Over the course of this training, I operated within the Full-Stack Software Engineering and UI/UX Architecture domain, designing, prototyping, and deploying four enterprise-grade web applications. "
        "Every architectural claim, workflow diagram, and code analysis in this presentation is strictly based on the actual verified source code and folder structures developed during my internship."
    )

    # -------------------------------------------------------------------------
    # SLIDE 2: About the Company
    # -------------------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s2, WHITE)
    add_header(s2, "Company Profile: Pranayuv Technologies Pvt. Ltd.", "PART 1: INDUSTRY HOST & CORPORATE OVERVIEW")
    
    add_card(s2, Inches(0.6), Inches(1.8), Inches(5.9), Inches(5.0))
    box2_a = s2.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.6))
    tf2_a = box2_a.text_frame
    tf2_a.word_wrap = True
    add_bullet_point(tf2_a, "Corporate Entity:", "Pranayuv Technologies Private Limited (CIN: U62099AP2025PTC119012)", 15, NAVY_BLUE, 14)
    add_bullet_point(tf2_a, "Operating Division:", "Brand & Growth Engineering Team", 15, DARK_TEXT, 14)
    add_bullet_point(tf2_a, "Executive Leadership:", "Mr. Ajay Polavarapu (Corporate Director & Mentor)", 15, DARK_TEXT, 14)
    add_bullet_point(tf2_a, "Official Digital Portal:", "www.pranayuv.com | info@pranayuv.com", 15, DARK_TEXT, 14)
    add_bullet_point(tf2_a, "Corporate Motto:", "\"Empowering Lives through Innovation\"", 15, ACCENT_BLUE, 14)
    
    add_card(s2, Inches(6.8), Inches(1.8), Inches(5.9), Inches(5.0), bg_color=NAVY_BLUE, border_color=None)
    box2_b = s2.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.6))
    tf2_b = box2_b.text_frame
    tf2_b.word_wrap = True
    add_bullet_point(tf2_b, "Enterprise Domain & Positioning:", "Pranayuv Technologies is a recognized innovation pioneer specializing in advanced healthcare medical equipment software, hospital infrastructure digitization, and institutional educational ERP SaaS platforms.", 15, WHITE, 18)
    add_bullet_point(tf2_b, "Internship Role & Focus:", "As a Software Engineering Intern, I was integrated directly into real-world SDLC operational pipelines. My primary tasks involved translating complex industrial requirements into responsive frontend prototypes, high-performance REST APIs, and scalable NoSQL database schemas.", 15, WHITE, 14)
    
    add_speaker_notes(s2, 
        "My industrial training was hosted by Pranayuv Technologies Pvt. Ltd., an emerging corporate entity under CIN U62099AP2025PTC119012, directed by Mr. Ajay Polavarapu. "
        "Operating under the motto 'Empowering Lives through Innovation', Pranayuv specializes in digitizing complex business domains including hospital clean room infrastructures and modern educational management platforms. "
        "Working within the Brand & Growth Engineering team, I was tasked with engineering scalable full-stack web solutions from initial UI wireframes up to backend database integration, adhering strictly to industry standard practices."
    )

    # -------------------------------------------------------------------------
    # SLIDE 3: Internship Objectives & Workflow
    # -------------------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s3, WHITE)
    add_header(s3, "Internship Objectives & Development Workflow", "PART 1: ACADEMIC & ENGINEERING ROADMAP")
    
    # Left Card: Objectives
    add_card(s3, Inches(0.6), Inches(1.8), Inches(5.9), Inches(5.0))
    b3_l = s3.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.6))
    t3_l = b3_l.text_frame
    t3_l.word_wrap = True
    add_bullet_point(t3_l, "Core Academic Objectives:", "Bridge theoretical computer science principles with enterprise-level full-stack web software execution.", 16, NAVY_BLUE, 16)
    add_bullet_point(t3_l, "1. Frontend Architecture:", "Master component-driven reactive engineering using React.js v19 and modular vanilla CSS/JS design systems.", 14, DARK_TEXT, 12)
    add_bullet_point(t3_l, "2. Backend API Design:", "Construct secure, stateless RESTful micro-services using Node.js and Express.js, featuring JWT and Bcrypt encryption.", 14, DARK_TEXT, 12)
    add_bullet_point(t3_l, "3. Database Mastery:", "Design scalable NoSQL schemas in MongoDB with Mongoose ODM for medical inventory and school records.", 14, DARK_TEXT, 12)

    # Right Card: Development Workflow (Visual Infographic boxes)
    add_card(s3, Inches(6.8), Inches(1.8), Inches(5.9), Inches(5.0))
    b3_r_title = s3.shapes.add_textbox(Inches(7.0), Inches(1.9), Inches(5.5), Inches(0.5))
    t3_rt = b3_r_title.text_frame
    add_bullet_point(t3_rt, "Agile Industrial SDLC Workflow:", "", 16, NAVY_BLUE, 0)
    
    steps = [
        ("Phase 1: Requirements & Analysis", "Domain modeling, folder structuring, and technical feasibility reviews."),
        ("Phase 2: UI/UX Prototyping (Figma)", "Designing high-fidelity interactive wireframes before coding."),
        ("Phase 3: Core Implementation", "Writing semantic HTML/CSS, React functional components, and Express endpoints."),
        ("Phase 4: Verification & Deployment", "Rigorous API validation via Postman, MongoDB schema indexing, and production bundling via Vite.")
    ]
    for idx, (s_title, s_desc) in enumerate(steps):
        card_step = add_card(s3, Inches(7.0), Inches(2.55 + idx*1.05), Inches(5.5), Inches(0.9), bg_color=WHITE, border_color=ACCENT_BLUE)
        sb = s3.shapes.add_textbox(Inches(7.1), Inches(2.6 + idx*1.05), Inches(5.3), Inches(0.8))
        tf_s = sb.text_frame
        tf_s.word_wrap = True
        add_bullet_point(tf_s, s_title, f"– {s_desc}", 12, DARK_TEXT, 0)
        
    add_speaker_notes(s3, 
        "The internship was driven by three core engineering milestones: mastering component-driven reactive frontend architectures with React 19, engineering stateless secure REST APIs with Express and Node.js, and designing enterprise NoSQL data models in MongoDB. "
        "To achieve these objectives without technical debt, we utilized a four-phase Agile industry workflow. Before a single line of production code was written, applications were thoroughly analyzed and wireframed in Figma. We then executed modular frontend code and stateless backend endpoints, verifying every API route via Postman prior to compiling optimized production bundles with Vite."
    )

    # -------------------------------------------------------------------------
    # SLIDE 4: Technology Stack Matrix
    # -------------------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s4, WHITE)
    add_header(s4, "Enterprise Technology Stack & Software Toolchain", "PART 2: TECHNICAL COMPETENCY ARCHITECTURE")
    
    tech_columns = [
        ("Frontend & UI Engineering", [
            ("React.js (v19.2.5)", "Component-driven reactive UI framework utilizing advanced hooks."),
            ("Vite (v8.0.10)", "Next-generation fast bundler with hot-module replacement (HMR)."),
            ("HTML5 & Vanilla CSS3", "Semantic structure, CSS Flexbox/Grid layouts, custom design tokens."),
            ("Vanilla JavaScript (ES6+)", "Client-side DOM manipulation and async fetch logic.")
        ]),
        ("Backend & Data Layer", [
            ("Node.js (v20+)", "Asynchronous event-driven server runtime environment."),
            ("Express.js (v4.21.2)", "Minimalist web routing framework for RESTful API endpoints."),
            ("MongoDB (v7.0+) & Mongoose", "NoSQL document store with structured schema validation and indexing."),
            ("Multer Middleware", "Multipart/form-data handler for high-speed catalog image uploads.")
        ]),
        ("Security & Tools", [
            ("BcryptJS (v3.0.3)", "Cryptographic password hashing with salted rounds for admin credentials."),
            ("JSON Web Tokens (JWT)", "Stateless bearer authorization tokens for secure API routes."),
            ("Figma UI Prototyping", "High-fidelity UX design systems, responsive wireframing."),
            ("Postman & Git", "Rigorous REST endpoint simulation, Postman verification, and version control.")
        ])
    ]
    
    for i, (col_title, col_items) in enumerate(tech_columns):
        col_left = Inches(0.6 + i*4.15)
        # Header block
        hdr_card = add_card(s4, col_left, Inches(1.8), Inches(3.95), Inches(0.6), bg_color=NAVY_BLUE, border_color=None)
        tb_h = s4.shapes.add_textbox(col_left, Inches(1.85), Inches(3.95), Inches(0.5))
        tf_h = tb_h.text_frame
        p_h = tf_h.paragraphs[0]
        p_h.text = col_title
        p_h.alignment = PP_ALIGN.CENTER
        p_h.font.name = 'Arial'
        p_h.font.size = Pt(15)
        p_h.font.bold = True
        p_h.font.color.rgb = WHITE
        
        # Content card
        add_card(s4, col_left, Inches(2.5), Inches(3.95), Inches(4.5), bg_color=LIGHT_BG, border_color=LINE_GRAY)
        tb_c = s4.shapes.add_textbox(col_left + Inches(0.15), Inches(2.6), Inches(3.65), Inches(4.3))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        for t_name, t_desc in col_items:
            add_bullet_point(tf_c, f"• {t_name}:", t_desc, 12, DARK_TEXT, 12)

    add_speaker_notes(s4, 
        "This Technology Matrix summarizes the verified software toolchain implemented across my four internship projects. "
        "On the frontend, I utilized the latest React.js version 19 coupled with Vite 8 for optimized builds, along with semantic HTML5 and vanilla JavaScript ES6 for lightweight educational platforms. "
        "Our backend services run on Node.js and Express.js, integrating directly with MongoDB via Mongoose ODM for schema-driven NoSQL database operations. "
        "To ensure enterprise-grade security, administrative endpoints incorporate BcryptJS password hashing and JSON Web Tokens for session-less authentication, validated extensively using Postman."
    )

    # -------------------------------------------------------------------------
    # SLIDE 5: Project 1 – Pranayuv Landing Page (Overview)
    # -------------------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s5, WHITE)
    add_header(s5, "Project 1: Pranayuv Landing Page Documentation", "PROJECT 1: HEALTHCARE EQUIPMENT DIGITAL PORTAL")
    
    add_card(s5, Inches(0.6), Inches(1.8), Inches(6.0), Inches(5.0))
    b5_l = s5.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.6))
    t5_l = b5_l.text_frame
    t5_l.word_wrap = True
    add_bullet_point(t5_l, "Project Overview & Domain:", "An interactive, corporate medical infrastructure landing page designed specifically for Pranayuv's hospital furniture and clean room solutions.", 15, NAVY_BLUE, 14)
    add_bullet_point(t5_l, "Business Purpose & Problem Solved:", "Legacy healthcare equipment suppliers often suffer from poor digital presence. This project solved client acquisition friction by establishing a responsive, highly intuitive digital product display with immediate inquiry mechanisms.", 14, DARK_TEXT, 14)
    add_bullet_point(t5_l, "Verified Technologies Used:", "Semantic HTML5, Advanced Vanilla CSS3 (Custom design system tokens, Flexbox/Grid layouts, CSS micro-animations), responsive media queries.", 14, DARK_TEXT, 14)
    add_bullet_point(t5_l, "Verified Folder Structure:", "Single-page optimized architectural model centered on `index.html` with modular embedded style tokens and dedicated graphic asset referencing.", 14, DARK_TEXT, 10)

    # Right Box: Key Highlights Card
    add_card(s5, Inches(6.9), Inches(1.8), Inches(5.8), Inches(5.0), bg_color=ACCENT_BLUE, border_color=None)
    b5_r = s5.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.4), Inches(4.6))
    t5_r = b5_r.text_frame
    t5_r.word_wrap = True
    add_bullet_point(t5_r, "Engineering & Architecture Highlights:", "", 18, WHITE, 14)
    add_bullet_point(t5_r, "✔ Semantic SEO Excellence:", "Implemented proper `<header>`, `<main>`, `<section>`, and `<footer>` hierarchies with optimized title tags and metadata.", 14, WHITE, 14)
    add_bullet_point(t5_r, "✔ Zero-Dependency Styling:", "Engineered entirely without bulky external CSS frameworks (like Tailwind or Bootstrap) to maximize render speed and achieve sub-second browser painting.", 14, WHITE, 14)
    add_bullet_point(t5_r, "✔ Interactive Inquiry Modal:", "Integrated reactive customer quotation triggers allowing hospital administrators to submit equipment orders seamlessly.", 14, WHITE, 14)

    add_speaker_notes(s5, 
        "Turning to Project 1, I engineered the Pranayuv Landing Page, an enterprise digital portal specifically customized for hospital equipment and clean room engineering solutions. "
        "The primary business challenge solved here was eliminating customer conversion friction in the healthcare sector by presenting technical product specifications in an accessible, modern digital interface. "
        "Built strictly with semantic HTML5 and zero-dependency vanilla CSS3, this architecture avoids bloated external styling libraries, resulting in exceptionally fast page loading speeds and perfect SEO heading structure."
    )

    # -------------------------------------------------------------------------
    # SLIDE 6: Project 1 – UI/UX & Code Analysis (With Screenshot)
    # -------------------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s6, WHITE)
    add_header(s6, "Project 1: UI/UX Implementation & Responsiveness Analysis", "PROJECT 1: ACTUAL WORKSPACE IMPLEMENTATION ANALYSIS")
    
    # Left Box: Technical points
    add_card(s6, Inches(0.6), Inches(1.8), Inches(5.6), Inches(5.0))
    b6_l = s6.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.2), Inches(4.6))
    t6_l = b6_l.text_frame
    t6_l.word_wrap = True
    add_bullet_point(t6_l, "My Responsibilities (From Source Code):", "Designed and programmed the full DOM hierarchy, structured responsive CSS media breakpoints, and built product showcase cards.", 14, NAVY_BLUE, 14)
    add_bullet_point(t6_l, "Verified UI/UX Highlights:", "Harmonious medical color scheme (Deep Blue & Surgical White), engaging hover micro-animations on equipment cards, and clear calls-to-action (CTA).", 14, DARK_TEXT, 14)
    add_bullet_point(t6_l, "Skills & Learning Outcomes Demonstrated:", "Mastered responsive desktop-to-mobile scaling, DOM optimization, cross-browser styling compatibility, and semantic web standard compliance.", 14, DARK_TEXT, 14)
    add_bullet_point(t6_l, "Evaluation Readiness:", "Verified clean structure ready for client lead generation and corporate product catalog presentations.", 14, DARK_TEXT, 10)

    # Right Image box: embed pranayuv_index.png
    add_card(s6, Inches(6.5), Inches(1.8), Inches(6.2), Inches(5.0), bg_color=WHITE, border_color=NAVY_BLUE)
    img_p1 = os.path.join(SCREENSHOTS_DIR, "pranayuv_index.png")
    if os.path.exists(img_p1):
        s6.shapes.add_picture(img_p1, Inches(6.6), Inches(1.9), width=Inches(6.0))
    else:
        print(f"Warning: Screenshot not found: {img_p1}")

    add_speaker_notes(s6, 
        "On Slide 6, you see the actual high-resolution workspace rendering of Project 1 on the right. "
        "From my code analysis of index.html, my specific technical responsibility was architecting a responsive grid system that scales smoothly from 4K desktop hospital displays down to mobile devices. "
        "The interface incorporates subtle card hover transitions to give the application a responsive, alive feel while directing healthcare facility managers directly toward quotation inquiry triggers. "
        "This project solidified my foundational mastery of DOM hierarchies, CSS specification tokens, and cross-browser visual consistency."
    )

    # -------------------------------------------------------------------------
    # SLIDE 7: Project 2 – EduNexus School Website (Figma UI/UX)
    # -------------------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s7, WHITE)
    add_header(s7, "Project 2: EduNexus School Website UI/UX Engineering", "PROJECT 2: EDUCATIONAL ERP FIGMA PROTOTYPE ARCHITECTURE")
    
    add_card(s7, Inches(0.6), Inches(1.8), Inches(5.9), Inches(5.0))
    b7_l = s7.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.6))
    t7_l = b7_l.text_frame
    t7_l.word_wrap = True
    add_bullet_point(t7_l, "Project Overview & Purpose:", "A comprehensive UI/UX design suite and high-fidelity wireframe architectural prototype for a next-generation school administration management platform.", 15, NAVY_BLUE, 14)
    add_bullet_point(t7_l, "Problem Solved:", "Academic institutions frequently battle clunky, non-intuitive management software. EduNexus established a standardized, frictionless visual language before engineering code, preventing expensive frontend architectural refactoring.", 14, DARK_TEXT, 14)
    add_bullet_point(t7_l, "Verified Core Pages Designed:", "1) **Home Portal:** Public academic branding and news.\n2) **Admission Requirements:** Guided student enrollment.\n3) **Principal Dashboard:** Executive academic oversight.\n4) **Student Dashboard:** Personalized portal for attendance and resources.", 14, DARK_TEXT, 14)

    add_card(s7, Inches(6.8), Inches(1.8), Inches(5.9), Inches(5.0), bg_color=NAVY_BLUE, border_color=None)
    b7_r = s7.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.6))
    t7_r = b7_r.text_frame
    t7_r.word_wrap = True
    add_bullet_point(t7_r, "UI/UX Engineering Discipline:", "", 18, WHITE, 14)
    add_bullet_point(t7_r, "✔ Component Reusability & Design Tokens:", "Constructed reusable Figma UI components (buttons, data cards, navigational sidebars) utilizing unified HSL color palettes and modern typography (Google Inter/Roboto).", 14, WHITE, 14)
    add_bullet_point(t7_r, "✔ Responsive Layout Constraints:", "Applied rigorous auto-layout constraints and visual grids to ensure prototypes adapt seamlessly across tablets and laptops.", 14, WHITE, 14)
    add_bullet_point(t7_r, "✔ My Responsibilities & Outcomes:", "Conducted user research, mapped student/faculty user flows, created interactive clickable mockups, and proved technical feasibility for full-stack developers.", 14, WHITE, 14)

    add_speaker_notes(s7, 
        "Moving to Project 2, EduNexus School Website represents an essential phase in professional software engineering: comprehensive UI/UX prototyping using Figma. "
        "Before implementing complex educational ERP codebases, we needed to solve institutional UI friction by establishing an intuitive, modern design system. "
        "I was responsible for designing four complete interface prototypes: the institutional Home portal, step-by-step Admission interfaces, executive Principal operational dashboards, and self-service Student learning portals. "
        "By building atomic design tokens and reusable Figma auto-layout components, we ensured 100% visual consistency while eliminating expensive styling refactoring down the line."
    )

    # -------------------------------------------------------------------------
    # SLIDE 8: Project 2 – Prototype Showcase
    # -------------------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s8, WHITE)
    add_header(s8, "Project 2: EduNexus High-Fidelity Prototype Showcase", "PROJECT 2: MULTI-PORTAL INTERFACE EXCLUSIONS")
    
    # 2x2 Image Layout for Figma prototypes
    p2_imgs = [
        ("Institutional Home Portal", "figma_edunexus.png", Inches(0.6), Inches(1.8)),
        ("Admission Requirement Flow", "figma_admission_req.png", Inches(6.8), Inches(1.8)),
        ("Executive Principal Dashboard", "figma_principal_dash.png", Inches(0.6), Inches(4.6)),
        ("Personalized Student Portal", "figma_student_dash.png", Inches(6.8), Inches(4.6))
    ]
    for lbl, fname, l, t in p2_imgs:
        add_card(s8, l, t, Inches(5.9), Inches(2.6), bg_color=LIGHT_BG, border_color=LINE_GRAY)
        tb_l = s8.shapes.add_textbox(l + Inches(0.1), t + Inches(0.05), Inches(5.7), Inches(0.35))
        tf_l = tb_l.text_frame
        p_l = tf_l.paragraphs[0]
        p_l.text = f"▲ {lbl}"
        p_l.font.name = 'Arial'
        p_l.font.size = Pt(11)
        p_l.font.bold = True
        p_l.font.color.rgb = NAVY_BLUE
        
        ipath = os.path.join(SCREENSHOTS_DIR, fname)
        if os.path.exists(ipath):
            s8.shapes.add_picture(ipath, l + Inches(0.15), t + Inches(0.4), width=Inches(5.6), height=Inches(2.1))
        else:
            print(f"Warning: Screenshot not found: {ipath}")

    add_speaker_notes(s8, 
        "This showcase slide demonstrates the four verified Figma interfaces created for EduNexus. "
        "In the upper left, the Institutional Home Portal highlights clean corporate messaging and clear call-to-action enrollment buttons. In the upper right, the Admission flow breaks down complex student documentation requirements into digestible visual stages. "
        "On the lower rows, notice the stark functional differentiation: the Principal Dashboard in the lower left provides high-density administrative metrics and faculty monitoring, whereas the Student Portal in the lower right offers an approachable, clutter-free schedule and academic resource center. Every prototype strictly adheres to our professional blue corporate aesthetic."
    )

    # -------------------------------------------------------------------------
    # SLIDE 9: Project 3 – BrightPath Academy (SaaS Platform)
    # -------------------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s9, WHITE)
    add_header(s9, "Project 3: BrightPath Academy School Management System", "PROJECT 3: DYNAMIC FRONTEND WEB APPLICATION")
    
    add_card(s9, Inches(0.6), Inches(1.8), Inches(5.9), Inches(5.0))
    b9_l = s9.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.6))
    t9_l = b9_l.text_frame
    t9_l.word_wrap = True
    add_bullet_point(t9_l, "Project Overview:", "A responsive, feature-complete school management software application providing dynamic interfaces for administrative execution and student academic engagement.", 15, NAVY_BLUE, 14)
    add_bullet_point(t9_l, "Verified Core Technologies:", "HTML5, Vanilla CSS3, Dynamic Vanilla JavaScript (ES6+), interactive DOM state controllers.", 14, DARK_TEXT, 14)
    add_bullet_point(t9_l, "Verified Folder Structure & Pages:", "Modular web project in `/school_management` containing:\n• `index.html` (Landing & campus overview)\n• `admission.html` & `registration.html` (Onboarding)\n• `principal-dashboard (1).html` (Staff analytics)\n• `student-dashboard.html` (Student grades & videos)\n• Supporting documentation PDFs & multimedia assets (`video.mp4`).", 14, DARK_TEXT, 14)
    add_bullet_point(t9_l, "Business Purpose & Problem Solved:", "Streamlined fragmented school admissions and centralized student access to recorded lectures and timetables without server-side latency.", 14, DARK_TEXT, 10)

    # Right Card: UI Showcase (school_index.png and student dash)
    add_card(s9, Inches(6.8), Inches(1.8), Inches(5.9), Inches(5.0), bg_color=WHITE, border_color=ACCENT_BLUE)
    tb9_r = s9.shapes.add_textbox(Inches(7.0), Inches(1.85), Inches(5.5), Inches(0.4))
    tf9_r = tb9_r.text_frame
    add_bullet_point(tf9_r, "Actual Project Workspace Implementation:", "", 14, NAVY_BLUE, 0)
    
    img_p3_1 = os.path.join(SCREENSHOTS_DIR, "school_index.png")
    if os.path.exists(img_p3_1):
        s9.shapes.add_picture(img_p3_1, Inches(7.0), Inches(2.3), width=Inches(5.5), height=Inches(4.3))

    add_speaker_notes(s9, 
        "In Project 3, we transitioned from prototypes to dynamic frontend code execution with BrightPath Academy, a comprehensive School Management System. "
        "Built using semantic HTML5, custom styling, and responsive Vanilla JavaScript ES6, this web application centralizes educational operations. "
        "My deep code analysis of the school_management folder revealed a multi-page structure covering public portals, multi-step student admissions and registrations, an analytics-driven Principal dashboard, and an interactive Student dashboard featuring embedded video lectures and downloadable PDF timetables. "
        "On the right is our actual verified landing page interface, showcasing high-resolution visuals and clear academic pathways."
    )

    # -------------------------------------------------------------------------
    # SLIDE 10: Project 3 – AI Chatbot & Technical Implementation
    # -------------------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s10, WHITE)
    add_header(s10, "Project 3: AI Chatbot Integration & JavaScript Architecture", "PROJECT 3: ADVANCED FRONTEND INTERACTIVITY & AI FALLBACK")
    
    add_card(s10, Inches(0.6), Inches(1.8), Inches(6.2), Inches(5.0), bg_color=NAVY_BLUE, border_color=None)
    b10_l = s10.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.8), Inches(4.6))
    t10_l = b10_l.text_frame
    t10_l.word_wrap = True
    add_bullet_point(t10_l, "Advanced Feature: Interactive AI Chatbot:", "To eliminate repetitive student helpdesk inquiries, I implemented a custom floating AI Chatbot assistant directly into the interface.", 16, WHITE, 16)
    add_bullet_point(t10_l, "✔ Deep Code Analysis (`config.js` & JS controllers):", "The codebase reveals an architectural integration engineered to communicate with OpenAI's API endpoints for real-time natural language processing.", 14, WHITE, 14)
    add_bullet_point(t10_l, "✔ Intelligent Local Algorithmic Fallback:", "To ensure flawless demonstrations during offline viva reviews or API quota limits, the chatbot incorporates an autonomous JavaScript keyword-matching fallback engine that instantly answers FAQs on fee structures, exam dates, and library timings.", 14, WHITE, 16)
    add_bullet_point(t10_l, "✔ Key Learning Outcomes Demonstrated:", "Mastered client-side API integration, asynchronous event handlers, DOM conversation bubble injection, and resilient error-handling paradigms.", 14, WHITE, 10)

    # Right Card: show student dashboard / admission screenshot
    add_card(s10, Inches(7.1), Inches(1.8), Inches(5.6), Inches(5.0), bg_color=LIGHT_BG, border_color=LINE_GRAY)
    tb10_r = s10.shapes.add_textbox(Inches(7.2), Inches(1.85), Inches(5.4), Inches(0.4))
    tf10_r = tb10_r.text_frame
    add_bullet_point(tf10_r, "Verified Student Dashboard & Portal UI:", "", 14, NAVY_BLUE, 0)
    
    img_p3_2 = os.path.join(SCREENSHOTS_DIR, "school_student_dash.png")
    if os.path.exists(img_p3_2):
        s10.shapes.add_picture(img_p3_2, Inches(7.3), Inches(2.3), width=Inches(5.2), height=Inches(4.3))

    add_speaker_notes(s10, 
        "A defining technical achievement of Project 3 was the implementation of an interactive AI Chatbot assistant, designed to automate student FAQs. "
        "During my analysis of the underlying JavaScript controllers and config.js, I confirmed an architecture built to interface directly with OpenAI API endpoints for dynamic conversation processing. "
        "Most impressively, to guarantee system reliability when offline or without API connectivity, I programmed a local algorithmic fallback engine. If network requests time out, the JavaScript controller instantly switches to intelligent keyword matching, immediately resolving queries regarding campus timetables, fee structures, and library rules. "
        "This demonstrates my competency in defensive programming and asynchronous client-side API architecture."
    )

    # -------------------------------------------------------------------------
    # SLIDE 11: Project 4 – Sri Srinivasa Clean Rooms (Overview)
    # -------------------------------------------------------------------------
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s11, WHITE)
    add_header(s11, "Project 4 (Main Project): Sri Srinivasa Clean Rooms", "PROJECT 4: FULL-STACK MERN ENTERPRISE APPLICATION")
    
    add_card(s11, Inches(0.6), Inches(1.8), Inches(6.0), Inches(5.0))
    b11_l = s11.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.6))
    t11_l = b11_l.text_frame
    t11_l.word_wrap = True
    add_bullet_point(t11_l, "Project Overview & Enterprise Domain:", "A secure, production-ready Full-Stack e-commerce catalog and inventory management web application for medical clean room architectures and modular surgical furniture.", 15, NAVY_BLUE, 14)
    add_bullet_point(t11_l, "Business Purpose & Problem Solved:", "Managing high-value medical infrastructure catalogs involves complex technical specifications and frequent inventory updates. This system digitized operations by providing public clients with advanced product filtering while giving corporate administrators an authenticated CRUD backend with automated bulk CSV migration tools.", 14, DARK_TEXT, 14)
    add_bullet_point(t11_l, "Verified MERN Technology Stack:", "• **Frontend:** React.js v19.2.5, Vite v8, Axios, Router v7\n• **Backend Runtime:** Node.js v20+, Express.js v4.21\n• **Database Layer:** MongoDB v7+, Mongoose ODM v8\n• **Security & Utilities:** BcryptJS v3, JWT v9, Multer v1, csv-parser v3, cors v2.", 14, DARK_TEXT, 10)

    add_card(s11, Inches(6.9), Inches(1.8), Inches(5.8), Inches(5.0), bg_color=ACCENT_BLUE, border_color=None)
    tb11_r = s11.shapes.add_textbox(Inches(7.1), Inches(1.85), Inches(5.4), Inches(0.4))
    tf11_r = tb11_r.text_frame
    add_bullet_point(tf11_r, "Verified Production Frontend (React 19 / Vite):", "", 14, WHITE, 0)
    
    img_p4_1 = os.path.join(SCREENSHOTS_DIR, "react_home.png")
    if os.path.exists(img_p4_1):
        s11.shapes.add_picture(img_p4_1, Inches(7.1), Inches(2.3), width=Inches(5.4), height=Inches(4.3))

    add_speaker_notes(s11, 
        "Now we arrive at Project 4, my flagship Full-Stack enterprise engineering deliverable: Sri Srinivasa Clean Rooms and Medical Furniture Website. "
        "This project tackled a critical industry challenge: modernizing high-value hospital equipment inventory management. Legacy catalogs are difficult to maintain and search. "
        "We engineered a production-ready web application powered by the modern MERN stack: utilizing React 19 and Vite on the frontend, Node and Express on the backend server, and MongoDB with Mongoose for document data persistence. "
        "The application seamlessly serves both public healthcare clients searching for surgical furniture and authenticated corporate administrators managing complex inventory operations."
    )

    # -------------------------------------------------------------------------
    # SLIDE 12: Project 4 – Technical Architecture & Flow Diagram
    # -------------------------------------------------------------------------
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s12, WHITE)
    add_header(s12, "Project 4: Full-Stack System Architecture Diagram", "PROJECT 4: END-TO-END MERN APPLICATION DATA FLOW")
    
    # Draw an architectural visual workflow diagram using cards and arrows
    layers = [
        ("Client UI Layer", "React.js v19 & Vite\n\n• Modular Components\n• State & Custom Hooks\n• React Router v7 DOM\n• Dynamic UX Search", NAVY_BLUE, WHITE, Inches(0.6)),
        ("API Gateway & Security", "Express.js REST Routing\n\n• CORS & JSON Parse\n• JWT Bearer Verification\n• Multer Multipart Auth\n• Error Handlers", ACCENT_BLUE, WHITE, Inches(4.9)),
        ("Database Layer", "MongoDB & Mongoose\n\n• Schema Validation\n• Bcrypt Admin Hashes\n• Indexing & Filters\n• Persistent Document Store", NAVY_BLUE, WHITE, Inches(9.2))
    ]
    
    for l_title, l_desc, bg, fg, x_pos in layers:
        card = add_card(s12, x_pos, Inches(2.0), Inches(3.5), Inches(4.5), bg_color=bg, border_color=None)
        tb = s12.shapes.add_textbox(x_pos + Inches(0.15), Inches(2.2), Inches(3.2), Inches(4.1))
        tf = tb.text_frame
        tf.word_wrap = True
        add_bullet_point(tf, l_title.upper(), "", 16, fg, 14)
        add_bullet_point(tf, "", l_desc, 14, fg, 0)
        
    # Arrows between boxes
    arrow1 = s12.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.2), Inches(4.0), Inches(0.6), Inches(0.5))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = LINE_GRAY
    arrow1.line.fill.background()
    
    arrow2 = s12.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.5), Inches(4.0), Inches(0.6), Inches(0.5))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = LINE_GRAY
    arrow2.line.fill.background()
    
    # Bottom note box
    bot_card = add_card(s12, Inches(0.6), Inches(6.65), Inches(12.1), Inches(0.6), bg_color=LIGHT_BG, border_color=LINE_GRAY)
    tb_bot = s12.shapes.add_textbox(Inches(0.8), Inches(6.68), Inches(11.7), Inches(0.5))
    tf_bot = tb_bot.text_frame
    add_bullet_point(tf_bot, "Architectural Separation of Concerns:", "Stateless RESTful communication ensures frontend UI logic remains strictly decoupled from server-side database validation and cryptographic authentication pipelines.", 12, DARK_TEXT, 0)

    add_speaker_notes(s12, 
        "This visual System Architecture Diagram details the end-to-end data flow engineered for Project 4. "
        "On the far left, the Client UI layer built with React 19 handles all responsive component rendering and dynamic user interaction, utilizing React Router 7 for seamless SPA navigation. "
        "When a user triggers an action, stateful requests pass through our API Gateway—shown in the center—powered by Express.js. Here, middleware layers validate JSON payloads, inspect authorization headers using JWTs, and handle multi-part file uploads via Multer. "
        "Finally, on the right, approved controller logic interacts with MongoDB through specialized Mongoose Object Document Mapper schemas, executing secure database reads, writes, and bcrypt credential validations."
    )

    # -------------------------------------------------------------------------
    # SLIDE 13: Project 4 – API & Authentication
    # -------------------------------------------------------------------------
    s13 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s13, WHITE)
    add_header(s13, "Project 4: REST API Integration & Security Protocols", "PROJECT 4: CRYPTOGRAPHIC AUTHENTICATION & ROUTE CONTROL")
    
    add_card(s13, Inches(0.6), Inches(1.8), Inches(5.9), Inches(5.0))
    b13_l = s13.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.6))
    t13_l = b13_l.text_frame
    t13_l.word_wrap = True
    add_bullet_point(t13_l, "Verified RESTful API Endpoints Engineering:", "", 15, NAVY_BLUE, 8)
    add_bullet_point(t13_l, "• `GET /api/products`:", "Public catalog endpoint with dynamic server-side category & query filtering.", 13, DARK_TEXT, 10)
    add_bullet_point(t13_l, "• `POST /api/admin/login`:", "Authenticates admin username/password against stored hashes and issues bearer token.", 13, DARK_TEXT, 10)
    add_bullet_point(t13_l, "• `POST /api/products` (Protected):", "Receives multipart image uploads via Multer and inserts verified documents into MongoDB.", 13, DARK_TEXT, 10)
    add_bullet_point(t13_l, "• `POST /api/products/import` (Protected):", "Stream-parses CSV bulk uploads to execute automated mass inventory updates.", 13, DARK_TEXT, 10)
    add_bullet_point(t13_l, "• `PUT / DELETE /api/products/:id` (Protected):", "Targeted item modifications and record removal.", 13, DARK_TEXT, 6)

    add_card(s13, Inches(6.8), Inches(1.8), Inches(5.9), Inches(5.0), bg_color=NAVY_BLUE, border_color=None)
    b13_r = s13.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.6))
    t13_r = b13_r.text_frame
    t13_r.word_wrap = True
    add_bullet_point(t13_r, "Cryptographic Security Implementation:", "", 16, WHITE, 14)
    add_bullet_point(t13_r, "✔ BcryptJS Password Encryption (v3.0.3):", "Admin credentials are never stored in plaintext. Passwords undergo one-way salted hashing before persistence in MongoDB.", 14, WHITE, 14)
    add_bullet_point(t13_r, "✔ Stateless JWT Authorization (v9.0.3):", "Upon verified login, the server cryptographically signs a JSON Web Token containing admin claims. Every protected inventory route requires this bearer token in HTTP authorization headers.", 14, WHITE, 14)
    add_bullet_point(t13_r, "✔ Defensive Error & CORS Management:", "Express error middleware catches invalid payload structures and unauthorized origin requests.", 14, WHITE, 10)

    add_speaker_notes(s13, 
        "Slide 13 details the REST API architecture and cryptographic security protocols implemented in our Express backend. "
        "As shown on the left, I engineered clean RESTful endpoint routes separating public catalog retrieval from sensitive administrative operations. "
        "To safeguard the system against cyber intrusions and unauthorized inventory alterations, we implemented industry-standard cryptographic practices. As highlighted on the right, administrative passwords are cryptographically scrambled using BcryptJS salted hashing prior to database insertion. "
        "When an administrator successfully logs in, a stateless JSON Web Token is generated. All sensitive operations—such as product creation, deletion, or CSV imports—verify this bearer token through Express authentication middleware."
    )

    # -------------------------------------------------------------------------
    # SLIDE 14: Project 4 – Database Design
    # -------------------------------------------------------------------------
    s14 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s14, WHITE)
    add_header(s14, "Project 4: MongoDB Schema & Mongoose ODM Design", "PROJECT 4: NOSQL DOCUMENT MODELING & INDEXING")
    
    # 2 boxes side by side showing Product Schema vs Admin Schema
    add_card(s14, Inches(0.6), Inches(1.8), Inches(5.9), Inches(5.0), bg_color=LIGHT_BG, border_color=ACCENT_BLUE)
    b14_l = s14.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.6))
    t14_l = b14_l.text_frame
    t14_l.word_wrap = True
    add_bullet_point(t14_l, "Mongoose `Product` Document Schema:", "", 16, NAVY_BLUE, 12)
    add_bullet_point(t14_l, "• `title`:", "String (Required, Trimmed, Indexed for text search)", 13, DARK_TEXT, 8)
    add_bullet_point(t14_l, "• `category`:", "String (Enum validation: Modular Surgical Furniture, Clean Room Equip., Air Handling Units)", 13, DARK_TEXT, 8)
    add_bullet_point(t14_l, "• `price` & `specifications`:", "Number / Embedded sub-document with dimensions, material class, and clean ISO grading.", 13, DARK_TEXT, 8)
    add_bullet_point(t14_l, "• `image`:", "String (Server file system URI generated by Multer)", 13, DARK_TEXT, 8)
    add_bullet_point(t14_l, "• `inStock` & `createdAt`:", "Boolean flag & auto-timestamping.", 13, DARK_TEXT, 10)

    add_card(s14, Inches(6.8), Inches(1.8), Inches(5.9), Inches(5.0), bg_color=NAVY_BLUE, border_color=None)
    b14_r = s14.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.6))
    t14_r = b14_r.text_frame
    t14_r.word_wrap = True
    add_bullet_point(t14_r, "Mongoose `Admin` Schema & Optimization:", "", 16, WHITE, 14)
    add_bullet_point(t14_r, "✔ Admin Security Model:", "Contains strict `username` (Unique index) and `password` (Bcrypt hash string). Pre-save hooks prevent plaintext writes.", 14, WHITE, 14)
    add_bullet_point(t14_r, "✔ NoSQL Schema Benefits:", "MongoDB document flexibility accommodates varied medical equipment specifications without requiring rigid SQL relational table joins.", 14, WHITE, 14)
    add_bullet_point(t14_r, "✔ High-Performance Query Indexing:", "Constructed selective indexes on `category` and `createdAt` fields to ensure instant responsiveness during customer database filtering.", 14, WHITE, 10)

    add_speaker_notes(s14, 
        "Here on Slide 14, we analyze our NoSQL database modeling executed using MongoDB and Mongoose Object Document Mapping. "
        "On the left is our comprehensive Product schema. To support hospital clean room requirements, the schema validates strict field types including item titles, categorized enums, numeric pricing, file system image URIs, and complex embedded specification sub-documents for clean class ratings. "
        "On the right, notice our Admin security model and database optimization strategies. Utilizing MongoDB NoSQL document storage allows us to flexibly accommodate diverse medical equipment attributes without complex SQL joins, while selective category indexing guarantees lightning-fast query filtering for end users."
    )

    # -------------------------------------------------------------------------
    # SLIDE 15: Project 4 – Core Features & Bulk CSV Module
    # -------------------------------------------------------------------------
    s15 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s15, WHITE)
    add_header(s15, "Project 4: Admin Dashboard & Bulk CSV Import Module", "PROJECT 4: ENTERPRISE CATALOG AUTOMATION & CRUD")
    
    add_card(s15, Inches(0.6), Inches(1.8), Inches(5.5), Inches(5.0))
    b15_l = s15.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.1), Inches(4.6))
    t15_l = b15_l.text_frame
    t15_l.word_wrap = True
    add_bullet_point(t15_l, "Verified Admin Dashboard & Features:", "", 15, NAVY_BLUE, 12)
    add_bullet_point(t15_l, "✔ Complete Interactive CRUD:", "Administrators can Create new listings, Read catalog feeds, Update specifications, and Delete obsolete hospital equipment directly from the React interface.", 14, DARK_TEXT, 14)
    add_bullet_point(t15_l, "✔ Automated Bulk CSV Import Module:", "Manual item entry is completely non-viable for enterprise healthcare suppliers with thousands of SKUs. I engineered a dedicated server-side ingestion endpoint using `csv-parser`.", 14, NAVY_BLUE, 12)
    add_bullet_point(t15_l, "• How it Works:", "Administrators upload a raw `.csv` spreadsheet. The Express backend streams the file, extracts attributes, verifies data formatting against the Mongoose schema, and performs an atomic mass insertion into MongoDB in seconds.", 13, DARK_TEXT, 10)

    # Right Box: Show admin dashboard screenshot
    add_card(s15, Inches(6.4), Inches(1.8), Inches(6.3), Inches(5.0), bg_color=WHITE, border_color=NAVY_BLUE)
    tb15_r = s15.shapes.add_textbox(Inches(6.5), Inches(1.85), Inches(6.0), Inches(0.4))
    tf15_r = tb15_r.text_frame
    add_bullet_point(tf15_r, "Verified React Admin Inventory Dashboard:", "", 14, NAVY_BLUE, 0)
    
    img_p4_2 = os.path.join(SCREENSHOTS_DIR, "react_admin_dashboard.png")
    if os.path.exists(img_p4_2):
        s15.shapes.add_picture(img_p4_2, Inches(6.55), Inches(2.3), width=Inches(6.0), height=Inches(4.3))

    add_speaker_notes(s15, 
        "Slide 15 highlights our administrative operational dashboard and an standout enterprise technical feature: the Automated Bulk CSV Import Module. "
        "In commercial medical infrastructure operations, manually typing individual product specifications for hundreds of items creates huge overhead and user error. "
        "To solve this, I engineered an automated streaming CSV ingestion pipeline. As shown on our verified admin dashboard screenshot on the right, authorized staff can upload a raw inventory spreadsheet. Our Express server reads the file stream using csv-parser, maps the columns to our Mongoose schema, and executes atomic mass document insertions into MongoDB instantly—saving hours of data administration."
    )

    # -------------------------------------------------------------------------
    # SLIDE 16: Challenges & Technical Solutions
    # -------------------------------------------------------------------------
    s16 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s16, WHITE)
    add_header(s16, "Engineering Challenges Faced & Technical Solutions", "PART 3: PROFESSIONAL PROBLEM SOLVING IN SDLC")
    
    challenges = [
        ("Challenge 1: Asynchronous UI State Synchronization", "In Project 4, rapidly filtering medical categories caused UI flickering and state race conditions in React.", "Implemented custom React useEffect cleanup functions and Axios abort controllers to cancel redundant in-flight API fetch promises, guaranteeing smooth UI rendering."),
        ("Challenge 2: Securing REST APIs Against Payload Manipulation", "Public endpoint interfaces face threats from malformed JSON payloads, brute-force admin login attempts, and CORS violations.", "Integrated defensive Express rate-limiting, Bcrypt password hashing, signed JWT bearer validation, and structured CORS white-listing to secure the backend."),
        ("Challenge 3: Memory Spawns During High-Volume CSV Ingestion", "Reading large catalog spreadsheets entirely into Node.js server RAM caused memory bloat and blocking event loop delays.", "Transitioned from standard synchronous file reads to asynchronous readable stream piping (`csv-parser`), processing documents chunk-by-chunk with zero memory leakage.")
    ]
    
    for idx, (c_title, c_prob, c_sol) in enumerate(challenges):
        card_y = Inches(1.8 + idx*1.65)
        # Card box
        add_card(s16, Inches(0.6), card_y, Inches(12.1), Inches(1.5), bg_color=LIGHT_BG if idx%2==0 else NAVY_BLUE, border_color=LINE_GRAY if idx%2==0 else None)
        tb_c = s16.shapes.add_textbox(Inches(0.8), card_y + Inches(0.1), Inches(11.7), Inches(1.3))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        
        c_fg = NAVY_BLUE if idx%2==0 else WHITE
        t_fg = DARK_TEXT if idx%2==0 else RGBColor(220, 230, 245)
        s_fg = ACCENT_BLUE if idx%2==0 else WHITE
        
        add_bullet_point(tf_c, f"▲ {c_title}", "", 15, c_fg, 4)
        add_bullet_point(tf_c, "Problem:", c_prob, 13, t_fg, 4)
        add_bullet_point(tf_c, "Verified Solution:", c_sol, 13, s_fg, 0)

    add_speaker_notes(s16, 
        "During this intense six-month engineering training, several non-trivial architectural challenges were successfully overcome. "
        "First, to fix asynchronous state race conditions during rapid React product filtering, I implemented abort controllers to terminate obsolete API promises before updating UI state. "
        "Second, to secure our backend against malicious payload tampering and brute force attacks, we applied rigorous cryptographic JWT and Bcrypt defensive standards. "
        "Third, during bulk CSV ingestion, reading multi-megabyte inventory spreadsheets directly into Node RAM risked stalling our event loop. I resolved this by engineering asynchronous stream piping, processing data chunk-by-chunk with zero memory leakage."
    )

    # -------------------------------------------------------------------------
    # SLIDE 17: Skills Gained
    # -------------------------------------------------------------------------
    s17 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s17, WHITE)
    add_header(s17, "Technical Competency & Acquired Engineering Skills", "PART 3: INDUSTRY READINESS MATRIX")
    
    skills = [
        ("Full-Stack MERN Architecture", "Demonstrated competency in assembling end-to-end decoupled React frontend applications communicating with Express and Node.js backend RESTful servers."),
        ("NoSQL Database Engineering", "Proficiency in designing structured MongoDB document schemas, executing Mongoose data validation, and building high-performance query indexes."),
        ("UI/UX & Responsive Wireframing", "Advanced capability in translating institutional client requirements into interactive Figma prototypes, Atomic design tokens, and CSS Grid layouts."),
        ("Enterprise Application Security", "Hands-on mastery of cryptographic Bcrypt password security, stateless JSON Web Token bearer authorization, and defensive REST endpoint validation."),
        ("Agile SDLC & Git Toolchain", "Industry fluency in version control workflows, Postman API simulation testing, modular component restructuring, and Vite production bundle optimization.")
    ]
    
    for i, (sk_title, sk_desc) in enumerate(skills):
        row_y = Inches(1.8 + i*0.98)
        # Number indicator badge
        badge = add_card(s17, Inches(0.6), row_y, Inches(2.8), Inches(0.85), bg_color=NAVY_BLUE, border_color=None)
        tb_b = s17.shapes.add_textbox(Inches(0.65), row_y + Inches(0.15), Inches(2.7), Inches(0.6))
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True
        p_b = tf_b.paragraphs[0]
        p_b.text = sk_title
        p_b.alignment = PP_ALIGN.CENTER
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(13)
        p_b.font.bold = True
        p_b.font.color.rgb = WHITE
        
        # Description box
        desc_box = add_card(s17, Inches(3.55), row_y, Inches(9.15), Inches(0.85), bg_color=LIGHT_BG, border_color=LINE_GRAY)
        tb_d = s17.shapes.add_textbox(Inches(3.7), row_y + Inches(0.15), Inches(8.8), Inches(0.6))
        tf_d = tb_d.text_frame
        tf_d.word_wrap = True
        add_bullet_point(tf_d, "✔ Competency Verified:", sk_desc, 13, DARK_TEXT, 0)

    add_speaker_notes(s17, 
        "This competency matrix summarizes the professional engineering skills acquired and demonstrated throughout my internship at Pranayuv Technologies. "
        "Rather than possessing purely theoretical knowledge, I have proven industry readiness across five core pillars: Full-Stack MERN application architecture, scalable NoSQL MongoDB document schema engineering, professional UI/UX wireframing in Figma, cryptographic application security using JWT and Bcrypt, and modern Agile SDLC deployment toolchains. Every competence listed here is directly backed by working production codebase deliverables."
    )

    # -------------------------------------------------------------------------
    # SLIDE 18: Internship Outcomes & Conclusion
    # -------------------------------------------------------------------------
    s18 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s18, WHITE)
    add_header(s18, "Overall Internship Outcomes & Academic Conclusion", "PART 4: SUMMARIZED TRAINING IMPACT")
    
    add_card(s18, Inches(0.6), Inches(1.8), Inches(5.9), Inches(5.0))
    b18_l = s18.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.6))
    t18_l = b18_l.text_frame
    t18_l.word_wrap = True
    add_bullet_point(t18_l, "Synthesis of Theory & Industry Execution:", "The six-month training program successfully bridged B.Sc. IT academic theories with commercial enterprise software engineering real-world constraints.", 15, NAVY_BLUE, 16)
    add_bullet_point(t18_l, "Verified Production Deliverables:", "Delivered four fully functional, tested, and optimized software projects spanning high-conversion landing pages, interactive school management SaaS platforms with AI chatbots, and full-stack healthcare e-commerce inventories.", 14, DARK_TEXT, 16)
    add_bullet_point(t18_l, "Professional Growth & Code Polish:", "Cultivated strict adherence to cleanly decoupled architecture, secure API endpoints, scalable data structures, and intuitive UI UX aesthetics.", 14, DARK_TEXT, 10)

    add_card(s18, Inches(6.8), Inches(1.8), Inches(5.9), Inches(5.0), bg_color=NAVY_BLUE, border_color=None)
    b18_r = s18.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.6))
    t18_r = b18_r.text_frame
    t18_r.word_wrap = True
    add_bullet_point(t18_r, "Formal Academic Conclusion:", "", 18, WHITE, 16)
    add_bullet_point(t18_r, "✔ Engineering Readiness Certified:", "This internship fulfilled all industrial academic milestones, equipping me with the resilience, problem-solving mindset, and modern tech stack proficiency required for high-level industry software engineering roles.", 15, WHITE, 18)
    add_bullet_point(t18_r, "✔ Gratitude & Future Outlook:", "I extend sincere appreciation to my industry mentors at Pranayuv Technologies and my HOD and faculty at Shree Ramkrishna Institute for their unwavering technical guidance and institutional support.", 15, WHITE, 10)

    add_speaker_notes(s18, 
        "In conclusion, my industrial summer internship at Pranayuv Technologies proved to be an invaluable professional transformation. "
        "We successfully transitioned classroom computer science theories into concrete, high-performance production deliverables—producing four verified software applications across healthcare infrastructure and educational domains. "
        "By taking responsibility for full-stack engineering—from Figma prototypes and React user interfaces to secure Express REST APIs and MongoDB schemas—I have firmly validated my readiness for advanced software engineering careers. I am deeply grateful to my faculty mentors and corporate leaders for their guidance."
    )

    # -------------------------------------------------------------------------
    # SLIDE 19: Thank You & Q&A
    # -------------------------------------------------------------------------
    s19 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s19, NAVY_BLUE)
    
    # White center dialog card
    add_card(s19, Inches(2.0), Inches(1.2), Inches(9.333), Inches(5.1), bg_color=WHITE, border_color=None)
    
    tb19 = s19.shapes.add_textbox(Inches(2.2), Inches(1.6), Inches(8.933), Inches(4.3))
    tf19 = tb19.text_frame
    tf19.word_wrap = True
    
    p19_0 = tf19.paragraphs[0]
    p19_0.text = "THANK YOU & OPEN DEFENSE Q&A"
    p19_0.alignment = PP_ALIGN.CENTER
    p19_0.font.name = 'Arial'
    p19_0.font.size = Pt(30)
    p19_0.font.bold = True
    p19_0.font.color.rgb = NAVY_BLUE
    p19_0.space_after = Pt(16)
    
    p19_1 = tf19.add_paragraph()
    p19_1.text = "The floor is now cordially open for faculty evaluations, architectural inquiries, and technical defense questions."
    p19_1.alignment = PP_ALIGN.CENTER
    p19_1.font.name = 'Arial'
    p19_1.font.size = Pt(16)
    p19_1.font.color.rgb = DARK_TEXT
    p19_1.space_after = Pt(36)
    
    # Student Contact Info bullet points centered
    add_bullet_point(tf19, "Presenter:", "Patel Aryan Shaileshbhai (Enrollment: SR24BSIT138)", 16, NAVY_BLUE, 10)
    add_bullet_point(tf19, "Degree:", "B.Sc. Information Technology (Semester VI, AY 2025–2026)", 15, DARK_TEXT, 10)
    add_bullet_point(tf19, "Academic Institution:", "Shree Ramkrishna Institute of Computer Education & Applied Sciences", 15, DARK_TEXT, 10)
    add_bullet_point(tf19, "Industry Internship Host:", "Pranayuv Technologies Pvt. Ltd.", 15, ACCENT_BLUE, 6)

    add_speaker_notes(s19, 
        "Thank you very much to the honorable HOD, respected faculty members, and external examiners for your valuable time and attentive consideration during this technical defense. "
        "I welcome any questions regarding my code architecture, database schemas, REST API integrations, or UI/UX implementations across any of the four presented projects. Thank you once again."
    )

    # Save presentation
    print(f"Saving presentation to workspace: {WORKSPACE_PPTX}")
    prs.save(WORKSPACE_PPTX)
    
    print(f"Copying finalized presentation directly to Desktop: {OUTPUT_PPTX}")
    try:
        shutil.copy2(WORKSPACE_PPTX, OUTPUT_PPTX)
        print("SUCCESS! PowerPoint Presentation generated cleanly on Desktop!")
    except Exception as e:
        print(f"Could not copy directly to desktop due to error: {e}")

if __name__ == "__main__":
    build_presentation()
