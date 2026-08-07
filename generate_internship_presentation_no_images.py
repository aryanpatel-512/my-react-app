import os, shutil
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Target Paths on Desktop
DESKTOP_DIR = r"C:\Users\Aryan\OneDrive\Desktop\Pranayuv"
OUTPUT_PPTX = os.path.join(DESKTOP_DIR, "Final_University_Internship_Presentation.pptx")
WORKSPACE_PPTX = os.path.abspath("Final_University_Internship_Presentation_NoImages.pptx")

# Corporate Color Palette (Blue & White Premium Corporate Theme)
NAVY_BLUE = RGBColor(31, 78, 121)     # #1F4E79 - Main Brand & Header
ACCENT_BLUE = RGBColor(46, 117, 182)  # #2E75B6 - Highlights & Subheadings
LIGHT_BG = RGBColor(242, 244, 248)    # #F2F4F8 - Card Backgrounds
WHITE = RGBColor(255, 255, 255)       # Pure White
DARK_TEXT = RGBColor(38, 38, 38)      # #262626 - Charcoal body text
LINE_GRAY = RGBColor(210, 215, 225)   # #D2D7E1 - Borders & Gridlines

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title_text, category_text="INTERNSHIP TECHNICAL DEFENSE & WORK DESCRIPTION"):
    # Category / Breadcrumb banner
    cat_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = 'Arial'
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_BLUE
    
    # Main Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.7), Inches(12), Inches(0.7))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(26)
    p_title.font.bold = True
    p_title.font.color.rgb = NAVY_BLUE
    
    # Decorative header rule line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY_BLUE
    line.line.color.rgb = NAVY_BLUE

def add_card(slide, left, top, width, height, bg_color=LIGHT_BG, border_color=LINE_GRAY):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
    else:
        card.line.fill.background()
    return card

def add_bullet_point(text_frame, bold_prefix, text_body, pt_size=14, color=DARK_TEXT, space_after=8):
    p = text_frame.add_paragraph() if len(text_frame.paragraphs[0].text) > 0 else text_frame.paragraphs[0]
    p.space_after = Pt(space_after)
    
    run_b = p.add_run()
    run_b.text = bold_prefix + (" " if bold_prefix else "")
    run_b.font.name = 'Arial'
    run_b.font.bold = True
    run_b.font.size = Pt(pt_size)
    run_b.font.color.rgb = color
    
    run_t = p.add_run()
    run_t.text = text_body
    run_t.font.name = 'Arial'
    run_t.font.bold = False
    run_t.font.size = Pt(pt_size)
    run_t.font.color.rgb = color

def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

def build_presentation():
    prs = Presentation()
    # Set widescreen dimensions 16:9 (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # -------------------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s1, NAVY_BLUE)
    
    bar1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(0.2), Inches(4.7))
    bar1.fill.solid()
    bar1.fill.fore_color.rgb = WHITE
    bar1.line.fill.background()
    
    tb1 = s1.shapes.add_textbox(Inches(1.3), Inches(1.4), Inches(11.5), Inches(4.7))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p0 = tf1.paragraphs[0]
    p0.text = "INDUSTRIAL INTERNSHIP TECHNICAL DEFENSE"
    p0.font.name = 'Arial'
    p0.font.size = Pt(32)
    p0.font.bold = True
    p0.font.color.rgb = WHITE
    p0.space_after = Pt(10)
    
    p1 = tf1.add_paragraph()
    p1.text = "Comprehensive Architectural Work Description & Project Analysis"
    p1.font.name = 'Arial'
    p1.font.size = Pt(22)
    p1.font.color.rgb = RGBColor(200, 220, 245)
    p1.space_after = Pt(38)
    
    p2 = tf1.add_paragraph()
    p2.text = "Presenter: Patel Aryan Shaileshbhai  |  Enrollment: SR24BSIT138"
    p2.font.name = 'Arial'
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.space_after = Pt(6)
    
    p3 = tf1.add_paragraph()
    p3.text = "Academic Program: B.Sc. Information Technology (Semester VI - AY 2025–2026)"
    p3.font.name = 'Arial'
    p3.font.size = Pt(16)
    p3.font.color.rgb = WHITE
    p3.space_after = Pt(6)
    
    p4 = tf1.add_paragraph()
    p4.text = "Institution: Shree Ramkrishna Institute of Computer Education & Applied Sciences, Sarvajanik University"
    p4.font.name = 'Arial'
    p4.font.size = Pt(15)
    p4.font.color.rgb = RGBColor(220, 230, 245)
    p4.space_after = Pt(6)
    
    p5 = tf1.add_paragraph()
    p5.text = "Host Corporate Organization: Pranayuv Technologies Pvt. Ltd."
    p5.font.name = 'Arial'
    p5.font.size = Pt(16)
    p5.font.bold = True
    p5.font.color.rgb = RGBColor(160, 205, 255)

    add_speaker_notes(s1, 
        "Good day respected HOD, Faculty Members, and External Examiners. My name is Patel Aryan Shaileshbhai, Enrollment number SR24BSIT138, currently in Semester VI of B.Sc. IT. "
        "Today, I present the formal technical defense and descriptive work breakdown of my six-month summer internship completed at Pranayuv Technologies Pvt. Ltd. "
        "In this presentation, I focus strictly on the engineering architecture, functional problem-solving, algorithm integration, and detailed work descriptions across four production web software deliverables developed during my training."
    )

    # -------------------------------------------------------------------------
    # SLIDE 2: About the Company
    # -------------------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s2, WHITE)
    add_header(s2, "Corporate Host: Pranayuv Technologies Pvt. Ltd.", "PART 1: INDUSTRY HOST & CORPORATE OVERVIEW")
    
    add_card(s2, Inches(0.6), Inches(1.8), Inches(5.9), Inches(5.0))
    b2_l = s2.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.6))
    t2_l = b2_l.text_frame
    t2_l.word_wrap = True
    add_bullet_point(t2_l, "Corporate Identity:", "Pranayuv Technologies Private Limited (CIN: U62099AP2025PTC119012)", 15, NAVY_BLUE, 14)
    add_bullet_point(t2_l, "Assigned Division:", "Brand & Growth Engineering Team", 15, DARK_TEXT, 14)
    add_bullet_point(t2_l, "Executive Supervision:", "Mr. Ajay Polavarapu (Corporate Director & Technical Mentor)", 15, DARK_TEXT, 14)
    add_bullet_point(t2_l, "Official Corporate Portal:", "www.pranayuv.com | careers@pranayuv.com", 15, DARK_TEXT, 14)
    add_bullet_point(t2_l, "Organizational Vision:", "\"Empowering Lives through Innovation\"", 15, ACCENT_BLUE, 14)
    
    add_card(s2, Inches(6.8), Inches(1.8), Inches(5.9), Inches(5.0), bg_color=NAVY_BLUE, border_color=None)
    b2_r = s2.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.6))
    t2_r = b2_r.text_frame
    t2_r.word_wrap = True
    add_bullet_point(t2_r, "Industry Positioning & Software Domains:", "Pranayuv Technologies is an established technology enterprise specializing in two critical industry sectors: medical hospital infrastructure software (clean room equipment platforms) and institutional educational ERP solutions.", 15, WHITE, 18)
    add_bullet_point(t2_r, "Internship Operational Scope:", "As a Full-Stack Engineering Intern, I was integrated into active software engineering pipelines. My core responsibility centered on architectural UI/UX wireframing, reactive frontend application programming, RESTful API engineering, and NoSQL database modeling.", 15, WHITE, 14)
    
    add_speaker_notes(s2, 
        "My industrial training was conducted at Pranayuv Technologies Pvt. Ltd., registered under CIN U62099AP2025PTC119012, under the guidance of Director Mr. Ajay Polavarapu. "
        "Guided by the corporate motto 'Empowering Lives through Innovation', Pranayuv engineers modern SaaS solutions tailored for healthcare clean room infrastructures and modern educational management institutions. "
        "Working within the Brand & Growth Engineering division, my technical mission was translating complex institutional client requirements into high-performance, maintainable software architectures."
    )

    # -------------------------------------------------------------------------
    # SLIDE 3: Internship Objectives
    # -------------------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s3, WHITE)
    add_header(s3, "Core Internship Objectives & Engineering Pillars", "PART 1: ACADEMIC MILESTONES & SKILL BUILDING")
    
    cols_s3 = [
        ("Pillar 1: Frontend UI/UX Engineering", [
            ("Reactive SPA Design:", "Master component-driven architectural patterns using React.js (v19) and Vite (v8)."),
            ("Design Systems:", "Construct custom vanilla CSS3 styling tokens without reliance on external bloated frameworks."),
            ("Responsive UX:", "Engineer adaptive interface layouts guaranteeing legibility across desktop, tablet, and mobile displays.")
        ]),
        ("Pillar 2: RESTful API & Server Logic", [
            ("Stateless Microservices:", "Design scalable, decoupled backend routing architectures using Node.js and Express.js."),
            ("Enterprise Security:", "Implement cryptographic Bcrypt password hashing and JSON Web Token (JWT) session-less authentication."),
            ("Stream Processing:", "Build automated server endpoints capable of real-time multi-part image and CSV document ingestion.")
        ]),
        ("Pillar 3: NoSQL Database Architecture", [
            ("Document Modeling:", "Design scalable MongoDB schemas utilizing Mongoose ODM for healthcare and educational domain models."),
            ("Query Optimization:", "Configure selective index generation on frequently filtered attributes to ensure sub-second search responses."),
            ("Data Integrity:", "Enforce strict schema field validation, enumeration restrictions, and atomic mass data operations.")
        ])
    ]
    
    for i, (col_title, items) in enumerate(cols_s3):
        col_left = Inches(0.6 + i*4.15)
        # Header box
        add_card(s3, col_left, Inches(1.8), Inches(3.95), Inches(0.7), bg_color=NAVY_BLUE, border_color=None)
        tb_h = s3.shapes.add_textbox(col_left, Inches(1.85), Inches(3.95), Inches(0.6))
        tf_h = tb_h.text_frame
        tf_h.word_wrap = True
        p_h = tf_h.paragraphs[0]
        p_h.text = col_title
        p_h.alignment = PP_ALIGN.CENTER
        p_h.font.name = 'Arial'
        p_h.font.size = Pt(15)
        p_h.font.bold = True
        p_h.font.color.rgb = WHITE
        
        # Body box
        add_card(s3, col_left, Inches(2.6), Inches(3.95), Inches(4.4), bg_color=LIGHT_BG, border_color=LINE_GRAY)
        tb_b = s3.shapes.add_textbox(col_left + Inches(0.15), Inches(2.7), Inches(3.65), Inches(4.2))
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True
        for b_title, b_desc in items:
            add_bullet_point(tf_b, f"✔ {b_title}", b_desc, 13, DARK_TEXT, 12)

    add_speaker_notes(s3, 
        "The primary academic objectives of this internship were structured around three distinct engineering pillars designed to transition theoretical computer science principles into practical software execution. "
        "Pillar 1 focused on Frontend UI/UX engineering, mastering declarative component architectures with React 19 and zero-dependency custom design systems. "
        "Pillar 2 centered on secure backend development, designing stateless Express REST APIs protected by Bcrypt encryption and JWT authorization. "
        "Pillar 3 emphasized NoSQL database engineering, constructing optimized MongoDB document schemas via Mongoose to enforce data integrity across high-volume healthcare and educational records."
    )

    # -------------------------------------------------------------------------
    # SLIDE 4: Development Workflow & SDLC
    # -------------------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s4, WHITE)
    add_header(s4, "Agile Software Development Lifecycle (SDLC) Workflow", "PART 1: INDUSTRIAL ENGINEERING METHODOLOGY")
    
    add_card(s4, Inches(0.6), Inches(1.8), Inches(12.1), Inches(0.8), bg_color=LIGHT_BG, border_color=ACCENT_BLUE)
    tb_intro = s4.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(11.7), Inches(0.6))
    tf_intro = tb_intro.text_frame
    tf_intro.word_wrap = True
    add_bullet_point(tf_intro, "Iterative Industrial Methodology:", "To deliver reliable software deliverables without technical debt or security vulnerabilities, all four internship projects adhered strictly to a 4-Phase Agile Software Engineering workflow.", 15, NAVY_BLUE, 0)
    
    # 4 horizontal phase boxes
    phases = [
        ("Phase 1: Analysis", "Requirements Analysis & Domain Schema Modeling:\n\nEvaluating institutional pain points, defining functional scope, structuring directories, and mapping Mongoose database schemas."),
        ("Phase 2: Prototyping", "UI/UX Figma Wireframing & Design Systems:\n\nConstructing atomic design tokens, interactive prototypes, and testing responsive interface behaviors before code implementation."),
        ("Phase 3: Implementation", "Full-Stack Code Execution (MERN & Modular DOM):\n\nProgramming semantic HTML/CSS, React functional components, stateful hooks, and Express RESTful backend endpoints."),
        ("Phase 4: Verification", "Rigorous Testing & Vite Production Deployment:\n\nVerifying endpoint JSON responses via Postman, testing JWT security boundaries, and bundling production code via Vite.")
    ]
    
    for i, (p_title, p_desc) in enumerate(phases):
        box_left = Inches(0.6 + i*3.1)
        # Header box
        add_card(s4, box_left, Inches(2.8), Inches(2.95), Inches(0.6), bg_color=NAVY_BLUE, border_color=None)
        tb_p = s4.shapes.add_textbox(box_left, Inches(2.88), Inches(2.95), Inches(0.5))
        tf_p = tb_p.text_frame
        p_p = tf_p.paragraphs[0]
        p_p.text = p_title.upper()
        p_p.alignment = PP_ALIGN.CENTER
        p_p.font.name = 'Arial'
        p_p.font.size = Pt(14)
        p_p.font.bold = True
        p_p.font.color.rgb = WHITE
        
        # Content card
        add_card(s4, box_left, Inches(3.5), Inches(2.95), Inches(3.5), bg_color=LIGHT_BG, border_color=LINE_GRAY)
        tb_c = s4.shapes.add_textbox(box_left + Inches(0.12), Inches(3.6), Inches(2.7), Inches(3.3))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        add_bullet_point(tf_c, "", p_desc, 13, DARK_TEXT, 0)
        
        # Draw connector arrow except for last
        if i < 3:
            arrow = s4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, box_left + Inches(3.0), Inches(4.9), Inches(0.08), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_BLUE
            arrow.line.fill.background()

    add_speaker_notes(s4, 
        "To ensure our software deliverables met professional production standards, we utilized a formalized four-phase Agile development lifecycle across every project. "
        "In Phase 1, we executed domain domain analysis and database modeling to scope institutional needs. In Phase 2, we created high-fidelity interactive Figma wireframes to eliminate UI friction before programming. "
        "In Phase 3, we wrote clean, modular production code across React frontends and Express backends. Finally, in Phase 4, every REST endpoint underwent comprehensive simulation testing via Postman prior to compiling optimized production bundles with Vite."
    )

    # -------------------------------------------------------------------------
    # SLIDE 5: Technology Stack & Toolchain Matrix
    # -------------------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s5, WHITE)
    add_header(s5, "Enterprise Technology Stack & Engineering Toolchain", "PART 2: VERIFIED TECHNICAL TOOL SET")
    
    tech_cols = [
        ("Core Frontend Architecture", [
            ("React.js (v19.2.5):", "Component-driven reactive UI framework utilizing functional declarative state hooks and modular composition."),
            ("Vite (v8.0.10):", "Next-generation fast bundler providing instantaneous Hot-Module Replacement (HMR) and optimized production minification."),
            ("HTML5 & CSS3:", "Semantic structural markup paired with zero-dependency custom HSL design tokens and Flexbox/Grid layouts."),
            ("Vanilla JavaScript (ES6+):", "Lightweight asynchronous client-side controllers and DOM manipulating state engines.")
        ]),
        ("Backend & Data Layer", [
            ("Node.js (v20+):", "Asynchronous event-driven server runtime providing scalable I/O execution for enterprise endpoints."),
            ("Express.js (v4.21.2):", "Minimalist web application routing framework structuring stateless REST API gateways and middleware pipelines."),
            ("MongoDB (v7.0+):", "High-performance NoSQL document data store accommodating flexible medical inventory specification schemas."),
            ("Mongoose ODM (v8.9.1):", "Object Document Mapper enforcing rigorous schema field validation, enumeration restrictions, and selective indexing.")
        ]),
        ("Security & Utilities", [
            ("BcryptJS (v3.0.3):", "Cryptographic library executing one-way salted password hashing to secure administrative credentials in MongoDB."),
            ("JSON Web Tokens (v9.0.3):", "Stateless authorization protocol signing cryptographic bearer tokens required for protected inventory operations."),
            ("Multer & CSV-Parser:", "Server-side middleware handling rapid multi-part image storage and streaming bulk spreadsheet ingestion."),
            ("Postman & Git Toolchain:", "Version-controlled codebase management coupled with automated REST API simulation testing.")
        ])
    ]
    
    for i, (t_title, t_items) in enumerate(tech_cols):
        col_left = Inches(0.6 + i*4.15)
        add_card(s5, col_left, Inches(1.8), Inches(3.95), Inches(0.6), bg_color=NAVY_BLUE, border_color=None)
        tb_t = s5.shapes.add_textbox(col_left, Inches(1.85), Inches(3.95), Inches(0.5))
        tf_t = tb_t.text_frame
        p_t = tf_t.paragraphs[0]
        p_t.text = t_title
        p_t.alignment = PP_ALIGN.CENTER
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = WHITE
        
        add_card(s5, col_left, Inches(2.5), Inches(3.95), Inches(4.5), bg_color=LIGHT_BG, border_color=LINE_GRAY)
        tb_b = s5.shapes.add_textbox(col_left + Inches(0.15), Inches(2.6), Inches(3.65), Inches(4.3))
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True
        for name, desc in t_items:
            add_bullet_point(tf_b, f"• {name}", desc, 12, DARK_TEXT, 10)

    add_speaker_notes(s5, 
        "Slide 5 articulates the comprehensive technology toolchain mastered during my internship. "
        "Our frontend software utilizes React version 19 and Vite 8 for reactive enterprise applications, supplemented by semantic HTML5 and ES6 JavaScript for lightweight portals. "
        "On the backend server side, we utilize Node.js and Express.js directly integrated with MongoDB and Mongoose ODM for scalable document storage and schema enforcement. "
        "To ensure strict defense against security exploits, administrative workflows implement BcryptJS salted password hashing, stateless JSON Web Tokens for bearer authentication, and streaming middleware for safe file ingestion."
    )

    # -------------------------------------------------------------------------
    # SLIDE 6: Project 1 – Pranayuv Landing Page (Overview & Purpose)
    # -------------------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s6, WHITE)
    add_header(s6, "Project 1: Pranayuv Landing Page Documentation", "PROJECT 1: HEALTHCARE EQUIPMENT DIGITAL PORTAL")
    
    add_card(s6, Inches(0.6), Inches(1.8), Inches(5.9), Inches(5.0))
    b6_l = s6.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.6))
    t6_l = b6_l.text_frame
    t6_l.word_wrap = True
    add_bullet_point(t6_l, "Project Identification & Domain:", "An interactive, corporate medical infrastructure landing page customized specifically for Pranayuv's hospital furniture and clean room equipment offerings.", 15, NAVY_BLUE, 14)
    add_bullet_point(t6_l, "Business Purpose & Problem Solved:", "Legacy healthcare equipment suppliers suffer from severe lead conversion friction due to traditional paper catalogs and unresponsive web presence. This project resolved client acquisition hurdles by designing a highly intuitive digital equipment catalog featuring immediate quotation inquiry mechanisms.", 14, DARK_TEXT, 14)
    add_bullet_point(t6_l, "Verified Folder & Asset Structure:", "Engineered around a single-page optimized architecture rooted at `/index.html`, utilizing embedded custom styling design tokens and streamlined local image asset referencing.", 14, DARK_TEXT, 10)

    add_card(s6, Inches(6.8), Inches(1.8), Inches(5.9), Inches(5.0), bg_color=NAVY_BLUE, border_color=None)
    b6_r = s6.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.6))
    t6_r = b6_r.text_frame
    t6_r.word_wrap = True
    add_bullet_point(t6_r, "Engineering Excellence & Value Proposition:", "", 18, WHITE, 14)
    add_bullet_point(t6_r, "✔ High-Speed Semantic Architecture:", "Constructed with clean HTML5 `<section>` hierarchies that maximize search engine indexing and ensure accessibility compliance across institutional browsers.", 14, WHITE, 14)
    add_bullet_point(t6_r, "✔ Zero External Library Dependency:", "Designed entirely without heavy styling frameworks (such as Bootstrap or Tailwind), completely eliminating unused utility bloat and achieving instantaneous browser DOM painting.", 14, WHITE, 14)
    add_bullet_point(t6_r, "✔ Direct Client Quotation Funnels:", "Integrated reactive interactive customer quotation triggers allowing hospital directors to select clean room specs and request vendor orders instantly.", 14, WHITE, 10)

    add_speaker_notes(s6, 
        "Moving into our specific software deliverables, Project 1 centered on engineering the Pranayuv Landing Page, an enterprise digital showcase tailored for healthcare clean room architectures. "
        "The core business challenge surmounted here was eliminating lead conversion delays inherent in traditional medical supply catalogs. We transformed the customer acquisition journey by creating a responsive digital portal with direct quotation triggers. "
        "Remarkably, by coding custom vanilla CSS3 without heavy third-party CSS frameworks, we eliminated bundle bloat completely, ensuring sub-second loading speeds on hospital terminal displays."
    )

    # -------------------------------------------------------------------------
    # SLIDE 7: Project 1 – Technical Work Description
    # -------------------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s7, WHITE)
    add_header(s7, "Project 1: Technical Work Description & DOM Analysis", "PROJECT 1: ACTUAL WORKSPACE ENGINEERING IMPLEMENTATION")
    
    add_card(s7, Inches(0.6), Inches(1.8), Inches(5.9), Inches(5.0))
    b7_l = s7.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.6))
    t7_l = b7_l.text_frame
    t7_l.word_wrap = True
    add_bullet_point(t7_l, "Documented Codebase Responsibilities:", "Engineered the full DOM layout hierarchy, designed custom HSL color palette variables, programmed responsive CSS media breakpoints, and styled product feature cards.", 15, NAVY_BLUE, 14)
    add_bullet_point(t7_l, "Responsive Layout & Media Queries:", "Implemented standard fluid design principles. Using custom `@media (max-width: 768px)` breakpoints, multi-column equipment showcases dynamically collapse into single-column vertical touch layouts without text distortion.", 14, DARK_TEXT, 14)
    add_bullet_point(t7_l, "Visual UI/UX Highlights:", "Harmonized a clinical color scheme composed of Surgical White backgrounds contrasted against Deep Corporate Navy text (`#1F4E79`), supported by engaging hover micro-animations on product selection boxes.", 14, DARK_TEXT, 10)

    add_card(s7, Inches(6.8), Inches(1.8), Inches(5.9), Inches(5.0), bg_color=LIGHT_BG, border_color=ACCENT_BLUE)
    b7_r = s7.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.6))
    t7_r = b7_r.text_frame
    t7_r.word_wrap = True
    add_bullet_point(t7_r, "Demonstrated Competencies & Outcomes:", "", 17, NAVY_BLUE, 14)
    add_bullet_point(t7_r, "✔ Cross-Browser Rendering Consistency:", "Verified structural uniformity across modern Chromium, FireFox, and Safari web browsing architectures.", 14, DARK_TEXT, 14)
    add_bullet_point(t7_r, "✔ Dom Tree Optimization & Accessibility:", "Maintained a shallow DOM tree depth to prevent layout reflow bottlenecks during dynamic scroll behavior, incorporating ARIA attributes and descriptive image alt-tags for screen readers.", 14, DARK_TEXT, 14)
    add_bullet_point(t7_r, "✔ Industry Readiness Verified:", "Proven capability to construct high-conversion corporate landing pages meeting commercial brand standards independently.", 14, DARK_TEXT, 10)

    add_speaker_notes(s7, 
        "Slide 7 provides a detailed engineering work description of my codebase contributions to Project 1. "
        "My specific programming responsibilities focused on architecting a fluid DOM hierarchy and implementing custom CSS media query breakpoints that adapt seamlessly from widescreen desktop hospital viewing down to mobile handheld devices. "
        "By enforcing a shallow DOM tree and utilizing hardware-accelerated CSS micro-animations, we achieved smooth visual feedback during product card hover states without incurring rendering lags. "
        "This deliverable verified my competency in cross-browser compatibility, W3C structural standards, and accessible UI engineering."
    )

    # -------------------------------------------------------------------------
    # SLIDE 8: Project 2 – EduNexus School Website (UI/UX Overview)
    # -------------------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s8, WHITE)
    add_header(s8, "Project 2: EduNexus School Website UI/UX Engineering", "PROJECT 2: EDUCATIONAL ERP FIGMA PROTOTYPE ARCHITECTURE")
    
    add_card(s8, Inches(0.6), Inches(1.8), Inches(5.9), Inches(5.0))
    b8_l = s8.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.6))
    t8_l = b8_l.text_frame
    t8_l.word_wrap = True
    add_bullet_point(t8_l, "Project Identification & Focus:", "A comprehensive UI/UX interface design suite and high-fidelity wireframe prototype architecture for a modern, next-generation school administration software platform.", 15, NAVY_BLUE, 14)
    add_bullet_point(t8_l, "Business Purpose & Problem Solved:", "Academic institutions frequently struggle with outdated, cluttered management software that frustrates faculty and students alike. EduNexus established a rigorous interface design phase before coding began, solving institutional UI friction and eliminating expensive frontend refactoring down the line.", 14, DARK_TEXT, 14)
    add_bullet_point(t8_l, "Design Tooling & Methodology:", "Created entirely inside **Figma** using atomized auto-layout constraints, unified HSL color tokens, and Google typography (Inter & Roboto).", 14, DARK_TEXT, 10)

    add_card(s8, Inches(6.8), Inches(1.8), Inches(5.9), Inches(5.0), bg_color=NAVY_BLUE, border_color=None)
    b8_r = s8.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.6))
    t8_r = b8_r.text_frame
    t8_r.word_wrap = True
    add_bullet_point(t8_r, "UI/UX Engineering Principles Applied:", "", 18, WHITE, 14)
    add_bullet_point(t8_r, "✔ Component-Driven Design Tokens:", "Engineered reusable Figma UI components (navigation sidebars, metric indicator cards, interactive form fields) that mirror React functional component architecture.", 14, WHITE, 14)
    add_bullet_point(t8_r, "✔ Responsive Layout Grids:", "Applied strict 12-column visual grid structures and fluid margin sizing to guarantee interface legibility across interactive whiteboards, tablets, and laptops.", 14, WHITE, 14)
    add_bullet_point(t8_r, "✔ Risk & Cost Mitigation:", "Validating user flows through high-fidelity prototypes allowed stakeholders to refine feature specifications without burning engineering hours on code rewrites.", 14, WHITE, 10)

    add_speaker_notes(s8, 
        "Moving to Project 2, EduNexus School Website showcases my capabilities in foundational UI/UX architectural wireframing using Figma. "
        "In enterprise software engineering, jumping directly into programming without standardized UI prototypes often leads to disjointed interfaces and costly codebase rewrites. "
        "EduNexus addressed institutional administrative friction by establishing a cohesive, intuitive visual language beforehand. Using atomized auto-layout components and strict HSL color tokens, we created interactive mockups that allowed stakeholders to test and validate user journeys prior to full-stack code implementation—saving significant development overhead."
    )

    # -------------------------------------------------------------------------
    # SLIDE 9: Project 2 – Work Description of Designed Portals
    # -------------------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s9, WHITE)
    add_header(s9, "Project 2: Work Description of Designed Educational Portals", "PROJECT 2: MULTI-PORTAL INTERFACE EXPLANATION & USER FLOWS")
    
    portals = [
        ("1. Institutional Home Portal", "Designed as the public digital identity of the educational organization. Features heroic value proposition branding, announcement ticker feeds, academic program showcases, and intuitive call-to-action buttons directing prospective candidates to admissions."),
        ("2. Admission Requirement Flow", "Engineered to demystify complex enrollment paperwork. Replaces confusing text blocks with step-by-step visual documentation checklists, fee structure transparency tables, and clear document submission guidelines to guide parents effortlessly."),
        ("3. Executive Principal Dashboard", "An administrative operational center designed for school leadership. Displays high-density academic KPI metrics, real-time faculty attendance tracking charts, tuition revenue summaries, and broadcast messaging controls in a clean, professional grid."),
        ("4. Personalized Student Portal", "A self-service academic hub customized for students. Eliminates visual clutter to present interactive timetable schedules, examination result report cards, assignment deadlines, and digital learning library access in an engaging layout.")
    ]
    
    for i, (p_title, p_desc) in enumerate(portals):
        col_x = Inches(0.6) if i % 2 == 0 else Inches(6.8)
        row_y = Inches(1.8) if i < 2 else Inches(4.6)
        
        card = add_card(s9, col_x, row_y, Inches(5.9), Inches(2.55), bg_color=LIGHT_BG if i%2==0 else NAVY_BLUE, border_color=LINE_GRAY if i%2==0 else None)
        tb_p = s9.shapes.add_textbox(col_x + Inches(0.15), row_y + Inches(0.12), Inches(5.6), Inches(2.3))
        tf_p = tb_p.text_frame
        tf_p.word_wrap = True
        
        t_color = NAVY_BLUE if i%2==0 else WHITE
        d_color = DARK_TEXT if i%2==0 else RGBColor(220, 230, 245)
        
        add_bullet_point(tf_p, f"▲ {p_title}", "", 15, t_color, 6)
        add_bullet_point(tf_p, "", p_desc, 13, d_color, 0)

    add_speaker_notes(s9, 
        "Slide 9 breaks down the functional work description across the four verified portal interfaces designed for EduNexus. "
        "In the upper left, the Institutional Home Portal structures brand storytelling and announcement broadcasts. In the upper right, the Admission Requirement flow transforms overwhelming enrollment documentation into clean, step-by-step visual checklists. "
        "On the bottom row, observe our functional differentiation: the Principal Dashboard in the lower left offers executive leadership high-density analytics on faculty attendance and institutional metrics, whereas the Student Portal in the lower right provides a clutter-free hub for examination results, timetables, and resource access."
    )

    # -------------------------------------------------------------------------
    # SLIDE 10: Project 3 – BrightPath Academy (SaaS Overview)
    # -------------------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s10, WHITE)
    add_header(s10, "Project 3: BrightPath Academy School Management System", "PROJECT 3: DYNAMIC FRONTEND WEB APPLICATION")
    
    add_card(s10, Inches(0.6), Inches(1.8), Inches(5.9), Inches(5.0))
    b10_l = s10.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.6))
    t10_l = b10_l.text_frame
    t10_l.word_wrap = True
    add_bullet_point(t10_l, "Project Identification & Domain:", "A dynamic, multi-page frontend school management software application providing interactive portals for institutional governance and student learning engagement.", 15, NAVY_BLUE, 14)
    add_bullet_point(t10_l, "Verified Tech Stack:", "Semantic HTML5, Custom Vanilla CSS3, Dynamic Vanilla JavaScript (ES6+) utilizing asynchronous Event Listeners and LocalStorage state controllers.", 14, DARK_TEXT, 14)
    add_bullet_point(t10_l, "Verified Codebase Repository (`/school_management`):", "• `index.html` (Institutional overview & entrance)\n• `admission.html` & `registration.html` (Onboarding flows)\n• `principal-dashboard (1).html` (Staff analytics)\n• `student-dashboard.html` (Academic resource access)\n• Direct multimedia embeds (`video.mp4`, `timetable.pdf`).", 14, DARK_TEXT, 14)

    add_card(s10, Inches(6.8), Inches(1.8), Inches(5.9), Inches(5.0), bg_color=ACCENT_BLUE, border_color=None)
    b10_r = s10.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.6))
    t10_r = b10_r.text_frame
    t10_r.word_wrap = True
    add_bullet_point(t10_r, "Business Value & Operational Impact:", "", 18, WHITE, 14)
    add_bullet_point(t10_r, "✔ Eliminating Paper-Based Admission Delays:", "Replaced slow manual enrollment processing with intelligent client-side validation forms, preventing incomplete data submission.", 14, WHITE, 16)
    add_bullet_point(t10_r, "✔ Centralized Digital Resource Distribution:", "Empowereed students with self-service streaming of video lecture modules and instant downloading of academic PDF syllabi without server latency.", 14, WHITE, 16)
    add_bullet_point(t10_r, "✔ Lightweight Browser Execution:", "Constructed purely on standard browser APIs, ensuring flawless performance on legacy computers and educational computer labs.", 14, WHITE, 10)

    add_speaker_notes(s10, 
        "In Project 3, we progressed from static design prototypes into full frontend web application execution with BrightPath Academy, a feature-complete School Management System. "
        "Built strictly with HTML5, CSS3, and dynamic Vanilla JavaScript ES6, this codebase completely digitizes administrative workflows. "
        "Our structural verification of the school_management repository revealed a comprehensive multi-page application spanning onboarding registration forms, staff analytical dashboards, and a feature-rich student learning portal complete with embedded recorded lecture videos and timetable PDFs. "
        "By utilizing native browser DOM APIs, this platform executes with zero server latency even on minimal lab computers."
    )

    # -------------------------------------------------------------------------
    # SLIDE 11: Project 3 – AI Chatbot & Technical Work Description
    # -------------------------------------------------------------------------
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s11, WHITE)
    add_header(s11, "Project 3: AI Chatbot Architecture & JS Implementation", "PROJECT 3: ADVANCED CLIENT-SIDE INTERACTIVITY & ALGORITHMIC FALLBACK")
    
    add_card(s11, Inches(0.6), Inches(1.8), Inches(5.9), Inches(5.0), bg_color=NAVY_BLUE, border_color=None)
    b11_l = s11.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.6))
    t11_l = b11_l.text_frame
    t11_l.word_wrap = True
    add_bullet_point(t11_l, "Standout Feature: Interactive AI Chatbot Module:", "To alleviate administrative burden caused by repetitive helpdesk queries, I integrated an automated conversational assistant directly into the interface.", 16, WHITE, 14)
    add_bullet_point(t11_l, "✔ Deep Code Analysis (`config.js` & Controllers):", "Inspection of the source codebase reveals an architectural configuration engineered to transmit student prompt payloads to OpenAI API endpoints for real-time natural language AI interpretation.", 14, WHITE, 16)
    add_bullet_point(t11_l, "✔ Autonomous Algorithmic Fallback Engine:", "To guarantee reliable execution during offline college viva reviews or API quota limits, I engineered a client-side JavaScript keyword-matching fallback algorithm. When API requests time out, the local engine seamlessly intercepts prompts to provide verified answers on fee structures, exam schedules, and library rules without error.", 14, WHITE, 10)

    add_card(s11, Inches(6.8), Inches(1.8), Inches(5.9), Inches(5.0), bg_color=LIGHT_BG, border_color=LINE_GRAY)
    b11_r = s11.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.6))
    t11_r = b11_r.text_frame
    t11_r.word_wrap = True
    add_bullet_point(t11_r, "Documented Frontend Work Responsibilities:", "", 16, NAVY_BLUE, 14)
    add_bullet_point(t11_r, "• Dynamic DOM State Manipulation:", "Programmed event listeners that inject interactive conversation dialogue bubbles into the UI dynamically without triggering page reloads.", 13, DARK_TEXT, 12)
    add_bullet_point(t11_r, "• Defensive Client-Side Form Validation:", "Constructed regex-driven script checkers on enrollment forms to verify student email formatting, contact numbers, and numeric graduation criteria before submission.", 13, DARK_TEXT, 12)
    add_bullet_point(t11_r, "• Acquired Engineering Takeaway:", "Proved foundational competence in integrating third-party APIs, handling asynchronous fetch promises, and building fault-tolerant fallback algorithms in modern ES6 JavaScript.", 13, DARK_TEXT, 10)

    add_speaker_notes(s11, 
        "A standout engineering achievement in Project 3 was the implementation of an interactive AI Chatbot designed to automate helpdesk operations. "
        "During my analysis of config.js and client scripts, I confirmed an architectural design engineered to interface directly with OpenAI API endpoints for real-time natural language conversation. "
        "Most importantly, to ensure fault tolerance when operating offline or during API network limits, I engineered an intelligent client-side algorithmic fallback engine. If network promises time out, our local JavaScript controller intercepts student prompts and performs algorithmic keyword matching, instantly delivering accurate responses regarding tuition fees, exam schedules, and library rules without throwing runtime exceptions."
    )

    # -------------------------------------------------------------------------
    # SLIDE 12: Project 4 – Sri Srinivasa Clean Rooms (MERN Overview)
    # -------------------------------------------------------------------------
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s12, WHITE)
    add_header(s12, "Project 4 (Main Project): Sri Srinivasa Clean Rooms", "PROJECT 4: FULL-STACK MERN ENTERPRISE WEB APPLICATION")
    
    add_card(s12, Inches(0.6), Inches(1.8), Inches(5.9), Inches(5.0))
    b12_l = s12.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.6))
    t12_l = b12_l.text_frame
    t12_l.word_wrap = True
    add_bullet_point(t12_l, "Project Identification & Domain Focus:", "A production-grade, secure Full-Stack e-commerce product catalog and inventory management platform tailored for hospital medical furniture and modular surgical clean room architectures.", 15, NAVY_BLUE, 14)
    add_bullet_point(t12_l, "Business Purpose & Solved Challenges:", "Managing high-value medical infrastructure equipment requires precise technical specification tracking and continuous inventory updates. This project modernized legacy operations by empowering healthcare purchasers with rapid dynamic product filtering while providing corporate administrators with a secure CRUD backend featuring automated spreadsheet data migration tools.", 14, DARK_TEXT, 14)
    add_bullet_point(t12_l, "Role-Based Access Governance (RBAC):", "Strictly isolates public catalog visualization from secure administrative inventory alteration endpoints.", 14, DARK_TEXT, 10)

    add_card(s12, Inches(6.8), Inches(1.8), Inches(5.9), Inches(5.0), bg_color=NAVY_BLUE, border_color=None)
    b12_r = s12.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.6))
    t12_r = b12_r.text_frame
    t12_r.word_wrap = True
    add_bullet_point(t12_r, "Verified Enterprise MERN Technology Stack:", "", 18, WHITE, 14)
    add_bullet_point(t12_r, "✔ Frontend Client (React 19 & Vite 8):", "Built with React functional components, modular hooks (`useState`, `useEffect`), Axios API communications, and React Router v7 DOM navigation.", 14, WHITE, 14)
    add_bullet_point(t12_r, "✔ Backend API Gateway (Node 20 & Express 4):", "Stateless asynchronous RESTful server routing handling JSON payloads and multipart data streams.", 14, WHITE, 14)
    add_bullet_point(t12_r, "✔ NoSQL Persistence & Security:", "MongoDB v7 database managed via Mongoose ODM v8, fortified by **BcryptJS (v3)** password encryption, **JWT (v9)** tokens, and **Multer** file upload middleware.", 14, WHITE, 10)

    add_speaker_notes(s12, 
        "We now advance to Project 4, my flagship Full-Stack enterprise engineering deliverable: the Sri Srinivasa Clean Rooms and Medical Furniture platform. "
        "This project tackled a major commercial engineering obstacle: managing complex, high-value hospital surgical equipment catalogs. Legacy operations struggle with tedious database updates and poor client search capabilities. "
        "To resolve this, we built a modern Full-Stack web application powered by the MERN stack—leveraging React 19 and Vite on the frontend, Node and Express on the backend server, and MongoDB with Mongoose for document persistence. "
        "The architecture enforces strict Role-Based Access Control, seamlessly separating public hospital client filtering from authenticated corporate administrative CRUD controls."
    )

    # -------------------------------------------------------------------------
    # SLIDE 13: Project 4 – Full-Stack System Architecture Diagram
    # -------------------------------------------------------------------------
    s13 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s13, WHITE)
    add_header(s13, "Project 4: Full-Stack System Architecture Diagram", "PROJECT 4: END-TO-END DECOUPLED DATA FLOW ARCHITECTURE")
    
    layers = [
        ("1. Client UI Layer", "React.js v19 & Vite\n\n• Single Page App (SPA)\n• Functional Modular Hooks\n• React Router v7 Navigation\n• Axios REST Client\n• Dynamic Category Filtering", NAVY_BLUE, WHITE, Inches(0.6)),
        ("2. API Gateway & Auth", "Express.js REST Server\n\n• Stateless API Routing\n• CORS & JSON Body Parsing\n• JWT Bearer Verification\n• Multer Multipart Storage\n• CSV-Parser Stream Pipeline", ACCENT_BLUE, WHITE, Inches(4.9)),
        ("3. Database Layer", "MongoDB & Mongoose\n\n• Document Data Persistence\n• Mongoose Schema Validation\n• Bcrypt Password Hashes\n• Selective Category Indexing\n• Atomic Mass CRUD Execution", NAVY_BLUE, WHITE, Inches(9.2))
    ]
    
    for l_title, l_desc, bg, fg, x_pos in layers:
        add_card(s13, x_pos, Inches(2.0), Inches(3.5), Inches(4.4), bg_color=bg, border_color=None)
        tb = s13.shapes.add_textbox(x_pos + Inches(0.15), Inches(2.2), Inches(3.2), Inches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        add_bullet_point(tf, l_title.upper(), "", 16, fg, 14)
        add_bullet_point(tf, "", l_desc, 14, fg, 0)
        
    # Draw custom arrows between blocks
    arrow1 = s13.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.2), Inches(3.9), Inches(0.6), Inches(0.5))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = LINE_GRAY
    arrow1.line.fill.background()
    
    arrow2 = s13.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.5), Inches(3.9), Inches(0.6), Inches(0.5))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = LINE_GRAY
    arrow2.line.fill.background()
    
    # Bottom architectural summary box
    add_card(s13, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.7), bg_color=LIGHT_BG, border_color=ACCENT_BLUE)
    tb_bot = s13.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6))
    tf_bot = tb_bot.text_frame
    tf_bot.word_wrap = True
    add_bullet_point(tf_bot, "Decoupled Separation of Concerns:", "Stateless HTTP RESTful JSON communication guarantees that frontend presentation UI logic remains totally isolated from server-side database verification and cryptographic authentication pipelines.", 13, DARK_TEXT, 0)

    add_speaker_notes(s13, 
        "This conceptual visual architectural block diagram illustrates the decoupled end-to-end data flow engineered for Project 4. "
        "On the far left, our Client UI layer developed in React 19 handles responsive component rendering, routing, and user interaction. When a client performs an action, Axios initiates asynchronous REST HTTP requests across our API Gateway—shown in the center block—running on Express.js. "
        "Here, custom middleware pipelines intercept requests to parse JSON bodies, verify JSON Web Tokens, or process file streams via Multer. "
        "Finally, validated requests pass into the far right block, interacting with MongoDB through specialized Mongoose ODM schemas to execute atomic database CRUD operations and Bcrypt hash comparisons."
    )

    # -------------------------------------------------------------------------
    # SLIDE 14: Project 4 – API & Security Work Description
    # -------------------------------------------------------------------------
    s14 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s14, WHITE)
    add_header(s14, "Project 4: REST API Routing & Cryptographic Security", "PROJECT 4: ENDPOINT ENGINEERING & DEFENSIVE PROTOCOLS")
    
    add_card(s14, Inches(0.6), Inches(1.8), Inches(5.9), Inches(5.0))
    b14_l = s14.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.6))
    t14_l = b14_l.text_frame
    t14_l.word_wrap = True
    add_bullet_point(t14_l, "Verified RESTful Server Endpoint Architecture:", "", 15, NAVY_BLUE, 10)
    add_bullet_point(t14_l, "• `GET /api/products`:", "Public catalog feed endpoint supporting dynamic server-side category and search term filtering.", 13, DARK_TEXT, 10)
    add_bullet_point(t14_l, "• `POST /api/admin/login`:", "Evaluates submitted credentials against stored bcrypt hashes; upon success, issues a signed JWT bearer token.", 13, DARK_TEXT, 10)
    add_bullet_point(t14_l, "• `POST /api/products` (Protected):", "Intercepts multipart image uploads via Multer, saves file system URIs, and commits verified documents into MongoDB.", 13, DARK_TEXT, 10)
    add_bullet_point(t14_l, "• `POST /api/products/import` (Protected):", "Streams uploaded CSV spreadsheets to execute automated mass catalog inventory insertions.", 13, DARK_TEXT, 10)
    add_bullet_point(t14_l, "• `PUT / DELETE /api/products/:id` (Protected):", "Executes targeted item parameter modifications or permanent record deletion.", 13, DARK_TEXT, 6)

    add_card(s14, Inches(6.8), Inches(1.8), Inches(5.9), Inches(5.0), bg_color=NAVY_BLUE, border_color=None)
    b14_r = s14.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.6))
    t14_r = b14_r.text_frame
    t14_r.word_wrap = True
    add_bullet_point(t14_r, "Cryptographic Application Security Protocols:", "", 16, WHITE, 14)
    add_bullet_point(t14_r, "✔ BcryptJS Password Encryption (v3.0.3):", "Administrative passwords are never written or evaluated in plaintext. Passwords undergo one-way salted algorithmic hashing prior to database storage, neutralizing database breach exposures.", 14, WHITE, 14)
    add_bullet_point(t14_r, "✔ Stateless JWT Authorization (v9.0.3):", "Session-less security architecture. Every administrative route utilizes Express middleware to intercept and cryptographically verify the JSON Web Token bearer header before granting resource modification access.", 14, WHITE, 14)
    add_bullet_point(t14_r, "✔ CORS & Defensive Middleware:", "Enforced Cross-Origin Resource Sharing restrictions and automated JSON parsing error trapping.", 14, WHITE, 10)

    add_speaker_notes(s14, 
        "Slide 14 details the REST API routing architecture and rigorous cryptographic security protocols engineered in our Express backend server. "
        "On the left, notice how endpoint routes strictly cleanly isolate public product catalog fetching from protected inventory modifications. "
        "To safeguard against unauthorized alterations and cybersecurity intrusions, we implemented robust defensive protocols shown on the right. Administrative credentials undergo one-way salted algorithmic encryption using BcryptJS before database insertion—meaning passwords never exist in plaintext. "
        "When an administrator authenticates via login, the server cryptographically signs a JSON Web Token. Every subsequent protected operation—such as adding items or performing CSV bulk imports—requires this verified token in HTTP headers."
    )

    # -------------------------------------------------------------------------
    # SLIDE 15: Project 4 – Database Design Work Description
    # -------------------------------------------------------------------------
    s15 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s15, WHITE)
    add_header(s15, "Project 4: MongoDB Schema & Mongoose ODM Design", "PROJECT 4: NOSQL DOCUMENT DATA MODELING & INDEXING")
    
    add_card(s15, Inches(0.6), Inches(1.8), Inches(5.9), Inches(5.0), bg_color=LIGHT_BG, border_color=ACCENT_BLUE)
    b15_l = s15.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.6))
    t15_l = b15_l.text_frame
    t15_l.word_wrap = True
    add_bullet_point(t15_l, "Mongoose `Product` Document Schema:", "", 16, NAVY_BLUE, 12)
    add_bullet_point(t15_l, "• `title`:", "String (Required, automatic whitespace trimming, custom indexed for high-speed text search).", 13, DARK_TEXT, 10)
    add_bullet_point(t15_l, "• `category`:", "String (Strict enumeration validation: Modular Surgical Furniture, Clean Room Equip., Air Handling Units).", 13, DARK_TEXT, 10)
    add_bullet_point(t15_l, "• `price` & `specifications`:", "Numeric values and embedded JSON sub-documents documenting physical dimensions, stainless steel material classes, and ISO clean room grading.", 13, DARK_TEXT, 10)
    add_bullet_point(t15_l, "• `image`:", "String (Local server file system path generated automatically by Multer file upload middleware).", 13, DARK_TEXT, 10)
    add_bullet_point(t15_l, "• `inStock` & `createdAt`:", "Boolean inventory availability flag and automatic epoch timestamping.", 13, DARK_TEXT, 10)

    add_card(s15, Inches(6.8), Inches(1.8), Inches(5.9), Inches(5.0), bg_color=NAVY_BLUE, border_color=None)
    b15_r = s15.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.6))
    t15_r = b15_r.text_frame
    t15_r.word_wrap = True
    add_bullet_point(t15_r, "Mongoose `Admin` Schema & NoSQL Optimization:", "", 16, WHITE, 14)
    add_bullet_point(t15_r, "✔ Admin Security Schema:", "Structures unique indexing on `username` and stores encrypted `password` hash strings, incorporating pre-save Mongoose validation middleware.", 14, WHITE, 14)
    add_bullet_point(t15_r, "✔ NoSQL Schema Flexibility:", "MongoDB document structure seamlessly accommodates varied medical furniture dimensions without requiring rigid SQL relational table joins.", 14, WHITE, 14)
    add_bullet_point(t15_r, "✔ Selective Database Indexing:", "Engineered compound B-Tree indexing across `category` and `createdAt` fields, ensuring instantaneous query evaluation during high-traffic catalog filtering.", 14, WHITE, 10)

    add_speaker_notes(s15, 
        "On Slide 15, we examine our NoSQL database modeling executed via MongoDB and Mongoose Object Document Mapping. "
        "On the left is our extensive Product schema. To support hospital equipment requirements, the schema enforces strict validation rules across item titles, categorized enums, numerical pricing, file system image references, and embedded specification sub-documents detailing ISO clean class ratings. "
        "On the right, notice our Admin credential schema and NoSQL optimization strategy. Utilizing flexible MongoDB document storage allows us to manage varied medical equipment attributes without complex SQL relational joins, while selective category indexing guarantees lightning-fast query filtering."
    )

    # -------------------------------------------------------------------------
    # SLIDE 16: Project 4 – Core Features & Bulk CSV Import Module
    # -------------------------------------------------------------------------
    s16 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s16, WHITE)
    add_header(s16, "Project 4: Admin Dashboard & Bulk CSV Import Module", "PROJECT 4: ENTERPRISE INVENTORY AUTOMATION & CRUD")
    
    add_card(s16, Inches(0.6), Inches(1.8), Inches(5.9), Inches(5.0))
    b16_l = s16.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.6))
    t16_l = b16_l.text_frame
    t16_l.word_wrap = True
    add_bullet_point(t16_l, "Comprehensive Interactive Inventory CRUD:", "", 16, NAVY_BLUE, 12)
    add_bullet_point(t16_l, "• Create:", "Administrators can inject new medical equipment listings with multi-part image attachments directly via responsive React dashboard modals.", 14, DARK_TEXT, 12)
    add_bullet_point(t16_l, "• Read & Filter:", "Real-time catalog synchronization displaying active inventory feeds with immediate filtering across medical equipment categories.", 14, DARK_TEXT, 12)
    add_bullet_point(t16_l, "• Update:", "Inline specification modifications allowing rapid alterations to pricing, stock availability flags, and clean room ISO grades.", 14, DARK_TEXT, 12)
    add_bullet_point(t16_l, "• Delete:", "Secure item retirement removing obsolete equipment listings from public client visualization immediately.", 14, DARK_TEXT, 10)

    add_card(s16, Inches(6.8), Inches(1.8), Inches(5.9), Inches(5.0), bg_color=LIGHT_BG, border_color=NAVY_BLUE)
    b16_r = s16.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.6))
    t16_r = b16_r.text_frame
    t16_r.word_wrap = True
    add_bullet_point(t16_r, "Standout Enterprise Feature: Bulk CSV Import Module:", "", 16, NAVY_BLUE, 12)
    add_bullet_point(t16_r, "✔ Resolving Manual Data Bottlenecks:", "Manually entering individual specifications for hundreds of hospital surgical SKUs creates immense administrative overhead and high human error rates. I engineered a streaming data ingestion endpoint using `csv-parser`.", 14, DARK_TEXT, 14)
    add_bullet_point(t16_r, "✔ Automated Stream Ingestion Pipeline:", "Authorized administrators upload a raw `.csv` inventory spreadsheet file through the dashboard. The Express server intercepts the file via Multer, asynchronously pipes the readable data stream through `csv-parser`, validates field formatting against Mongoose schema rules, and executes atomic mass document insertions into MongoDB in seconds.", 13, DARK_TEXT, 10)

    add_speaker_notes(s16, 
        "Slide 16 details our interactive administrative CRUD dashboard and an advanced technical capability: the Automated Bulk CSV Import Module. "
        "On the left, notice how the React administrative interface equips managers with full CRUD operations over product listings and multi-part image uploads. "
        "However, in commercial healthcare supply chain operations, manually typing product data for hundreds of inventory SKUs causes massive bottlenecks and high user error rates. "
        "To resolve this, I engineered an automated streaming CSV ingestion endpoint. Authorized staff upload a raw inventory spreadsheet; our Express backend streams the file using csv-parser, maps columns directly to our Mongoose schema, and executes atomic mass document insertions into MongoDB in seconds—reducing data management overhead from hours to moments."
    )

    # -------------------------------------------------------------------------
    # SLIDE 17: Engineering Challenges & Technical Solutions
    # -------------------------------------------------------------------------
    s17 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s17, WHITE)
    add_header(s17, "Engineering Challenges Surmounted & Technical Solutions", "PART 3: PROFESSIONAL PROBLEM SOLVING IN SDLC")
    
    challenges_data = [
        ("Challenge 1: Asynchronous State Race Conditions in React", "Rapidly switching product filter tabs in Project 4 triggered out-of-order asynchronous API responses, causing severe UI flickering and invalid inventory catalog display.", "Implemented clean React `useEffect` abort controllers and cancellation cleanup routines. When a user changes filters, redundant in-flight API fetch promises are canceled instantly, guaranteeing deterministic UI state rendering."),
        ("Challenge 2: Hardening REST Endpoints Against Payload Attacks", "Exposed public backend routes faced severe cyber exposure from malformed JSON injection payloads, brute-force administrative credential attacks, and unauthorized origin CORS requests.", "Fortified Express gateway pipelines with structured rate-limiting, Bcrypt password encryption, stateless signed JWT bearer authorization validation, and restrictive origin CORS whitelisting."),
        ("Challenge 3: Server RAM Bloat During High-Volume CSV Ingestion", "Reading massive catalog spreadsheets synchronously into Node.js server RAM during bulk inventory uploads threatened memory exhaustion and blocking event loop latency.", "Engineered asynchronous readable file stream piping using `csv-parser`. Spreadsheets are parsed and pushed to MongoDB in small sequential data chunks, achieving high-speed processing with zero memory leakage.")
    ]
    
    for idx, (title, prob, sol) in enumerate(challenges_data):
        row_y = Inches(1.8 + idx*1.65)
        add_card(s17, Inches(0.6), row_y, Inches(12.1), Inches(1.5), bg_color=LIGHT_BG if idx%2==0 else NAVY_BLUE, border_color=LINE_GRAY if idx%2==0 else None)
        tb_c = s17.shapes.add_textbox(Inches(0.8), row_y + Inches(0.08), Inches(11.7), Inches(1.35))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        
        t_fg = NAVY_BLUE if idx%2==0 else WHITE
        p_fg = DARK_TEXT if idx%2==0 else RGBColor(220, 230, 245)
        s_fg = ACCENT_BLUE if idx%2==0 else WHITE
        
        add_bullet_point(tf_c, f"▲ {title}", "", 15, t_fg, 4)
        add_bullet_point(tf_c, "Technical Hurdle:", prob, 13, p_fg, 4)
        add_bullet_point(tf_c, "Engineered Solution:", sol, 13, s_fg, 0)

    add_speaker_notes(s17, 
        "During this intense six-month software engineering training, several non-trivial architectural hurdles were encountered and systematically conquered. "
        "First, to solve asynchronous state race conditions during rapid React category filtering, I implemented promise abort controllers that automatically cancel outdated network fetch calls before updating UI state. "
        "Second, to secure our endpoints against payload tampering and brute-force attacks, we integrated rigorous Bcrypt password hashing and JWT defensive authorization barriers. "
        "Third, during bulk CSV file uploads, loading entire spreadsheets directly into server RAM risked stalling our Node event loop. I eliminated this memory bloat by engineering asynchronous stream piping, processing records chunk-by-chunk with zero memory leakage."
    )

    # -------------------------------------------------------------------------
    # SLIDE 18: Acquired Competencies & Internship Outcomes
    # -------------------------------------------------------------------------
    s18 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s18, WHITE)
    add_header(s18, "Acquired Engineering Competencies & Academic Outcomes", "PART 3 & 4: CERTIFIED INDUSTRY READINESS & IMPACT")
    
    add_card(s18, Inches(0.6), Inches(1.8), Inches(5.9), Inches(5.0))
    b18_l = s18.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.6))
    t18_l = b18_l.text_frame
    t18_l.word_wrap = True
    add_bullet_point(t18_l, "Verified Professional Technical Competencies:", "", 16, NAVY_BLUE, 12)
    add_bullet_point(t18_l, "• Full-Stack MERN Architecture:", "Proven mastery in structuring decoupled React single-page frontends integrated with Node/Express REST backend routing gateways.", 13, DARK_TEXT, 10)
    add_bullet_point(t18_l, "• NoSQL Schema Engineering:", "Proficiency in modeling MongoDB document databases, setting selective indexes, and enforcing Mongoose data validation.", 13, DARK_TEXT, 10)
    add_bullet_point(t18_l, "• UI/UX & Responsive Wireframing:", "Advanced capability translating complex client requirements into intuitive Figma design systems and fluid CSS Grid layouts.", 13, DARK_TEXT, 10)
    add_bullet_point(t18_l, "• Cryptographic Application Security:", "Hands-on competency in executing Bcrypt salted hashing, stateless JWT authorization, and defensive input validation.", 13, DARK_TEXT, 10)
    add_bullet_point(t18_l, "• Agile SDLC & Version Control:", "Professional fluency in Git version management, Postman simulation validation, and Vite deployment toolchains.", 13, DARK_TEXT, 6)

    add_card(s18, Inches(6.8), Inches(1.8), Inches(5.9), Inches(5.0), bg_color=NAVY_BLUE, border_color=None)
    b18_r = s18.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.6))
    t18_r = b18_r.text_frame
    t18_r.word_wrap = True
    add_bullet_point(t18_r, "Summarized Academic & Industrial Impact:", "", 17, WHITE, 14)
    add_bullet_point(t18_r, "✔ Theory-to-Execution Bridge:", "This internship successfully transformed classroom B.Sc. IT foundational computer science principles into commercial enterprise software engineering capabilities.", 14, WHITE, 16)
    add_bullet_point(t18_r, "✔ Concrete Software Portfolio:", "Delivered four fully verified, high-performance production deliverables across healthcare infrastructure and educational SaaS platforms.", 14, WHITE, 16)
    add_bullet_point(t18_r, "✔ Formal Certification of Readiness:", "Fully prepared to assume advanced software engineering industry responsibilities with rigorous architectural discipline and problem-solving resilience.", 14, WHITE, 10)

    add_speaker_notes(s18, 
        "Slide 18 synthesizes the technical competencies acquired and the formal academic outcomes achieved throughout my internship at Pranayuv Technologies. "
        "Rather than possessing purely abstract classroom knowledge, I have proven industry readiness across five foundational pillars: Full-Stack MERN application architecture, scalable NoSQL MongoDB schema engineering, UI/UX wireframing in Figma, cryptographic API security, and Agile Git toolchains. "
        "By delivering four working production codebases across healthcare and educational domains, this six-month training has thoroughly validated my preparedness for high-level industry software engineering careers."
    )

    # -------------------------------------------------------------------------
    # SLIDE 19: Thank You & Q&A
    # -------------------------------------------------------------------------
    s19 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s19, NAVY_BLUE)
    
    add_card(s19, Inches(2.0), Inches(1.2), Inches(9.333), Inches(5.1), bg_color=WHITE, border_color=None)
    tb19 = s19.shapes.add_textbox(Inches(2.2), Inches(1.6), Inches(8.933), Inches(4.3))
    tf19 = tb19.text_frame
    tf19.word_wrap = True
    
    p19_0 = tf19.paragraphs[0]
    p19_0.text = "THANK YOU & OPEN DEFENSE Q&A"
    p19_0.alignment = PP_ALIGN.CENTER
    p19_0.font.name = 'Arial'
    p19_0.font.size = Pt(30)
    p19_0.font.bold = True
    p19_0.font.color.rgb = NAVY_BLUE
    p19_0.space_after = Pt(16)
    
    p19_1 = tf19.add_paragraph()
    p19_1.text = "The floor is now cordially open for faculty evaluations, architectural inquiries, and technical defense questions."
    p19_1.alignment = PP_ALIGN.CENTER
    p19_1.font.name = 'Arial'
    p19_1.font.size = Pt(16)
    p19_1.font.color.rgb = DARK_TEXT
    p19_1.space_after = Pt(36)
    
    add_bullet_point(tf19, "Presenter:", "Patel Aryan Shaileshbhai (Enrollment: SR24BSIT138)", 16, NAVY_BLUE, 10)
    add_bullet_point(tf19, "Academic Degree:", "B.Sc. Information Technology (Semester VI, AY 2025–2026)", 15, DARK_TEXT, 10)
    add_bullet_point(tf19, "University & Institution:", "Shree Ramkrishna Institute of Computer Education & Applied Sciences, Sarvajanik University", 15, DARK_TEXT, 10)
    add_bullet_point(tf19, "Hosting Corporate Organization:", "Pranayuv Technologies Pvt. Ltd.", 15, ACCENT_BLUE, 6)

    add_speaker_notes(s19, 
        "Thank you very much to the honorable HOD, respected faculty members, and external examiners for your valuable time and attentive consideration during this technical defense. "
        "I cordially welcome any evaluations or questions regarding my code architecture, database schemas, REST API security implementations, or algorithmic fallback logic across any of the four presented software deliverables. Thank you once again."
    )

    # Save presentation
    print(f"Saving image-free technical presentation to workspace: {WORKSPACE_PPTX}")
    prs.save(WORKSPACE_PPTX)
    
    print(f"Copying finalized presentation directly to Desktop: {OUTPUT_PPTX}")
    try:
        shutil.copy2(WORKSPACE_PPTX, OUTPUT_PPTX)
        print("SUCCESS! PowerPoint Presentation generated cleanly on Desktop!")
    except Exception as e:
        print(f"Could not copy directly to desktop due to error: {e}")

if __name__ == "__main__":
    build_presentation()
